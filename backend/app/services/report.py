"""
报告服务

负责报告的 CRUD 及多格式导出（Markdown → DOCX/PDF）。

导出功能说明：
- generate_docx(): 使用 python-docx 生成 Word 文档，按 Markdown 标题层级映射样式
- generate_pdf(): 使用 reportlab 生成 PDF，标题自动映射到对应字体样式
- generate_markdown(): 直接拼接 Markdown 内容返回原始文本
- export_report(): 统一入口，根据 fmt 参数路由到对应生成器

为什么同时支持 DOCX 和 PDF：
- DOCX 面向编辑（用户下载后可在 Word 中继续修改）
- PDF 面向分发（格式固定，跨设备查看体验一致）
- Markdown 面向集成（可直接导入 Notion、Obsidian 等知识管理工具）
"""

import asyncio
import base64
import io
import logging
import os
import re
import tempfile
from typing import Any, Dict, List, Optional, Tuple
import docx
import docx.shared
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from sqlalchemy import select, or_, func
from sqlalchemy.ext.asyncio import AsyncSession
from app.models.report import Report
from app.schemas.report import ReportCreate, ReportUpdate, ReportStatisticsItem
from app.services.ppt_template import PPTTemplateService
from app.services.report_page_model import (
    ReportPageModel, PageModel, TextBlock, ImageBlock,
    USABLE_WIDTH, USABLE_HEIGHT,
)
from app.services.embedding import EmbeddingService
from app.services.vectorstore import VectorStoreService

logger = logging.getLogger(__name__)

class ReportService:
    @staticmethod
    def _chunk_text(text: str, chunk_size: int = 512, overlap: int = 64) -> list[str]:
        if not text:
            return []
        chunks = []
        start = 0
        while start < len(text):
            end = min(start + chunk_size, len(text))
            chunks.append(text[start:end])
            start += chunk_size - overlap
        return chunks

    async def _embed_and_store(self, report_id: int, content_markdown: str):
        """DEPRECATED: 旧版自动嵌入入口，由 export_report 触发。

        自"lazy 向量库写入"改造后，写入时机已统一迁移到
        `VectorstoreUploadService.upload_report()`，并由 export / upload 端点
        在用户主动确认时触发。本方法保留仅为向后兼容，新逻辑请走
        `upload_to_vector_store_if_pending()`。
        """
        try:
            chunks = self._chunk_text(content_markdown)
            if not chunks:
                return
            emb_service = EmbeddingService()
            vectors = await emb_service.embed_texts_or_empty(chunks)
            if not vectors:
                logger.warning("Embedding returned empty, skipping vector store for report %d", report_id)
                return
            vs_service = VectorStoreService()
            if not vs_service._connected:
                logger.warning("Milvus not connected, skipping embedding for report %d", report_id)
                return
            vs_service.init_collection(len(vectors[0]))
            vs_service.insert(str(report_id), chunks, vectors)
            logger.info("Embedded %d chunks for report %d", len(chunks), report_id)
        except Exception as e:
            logger.warning("Failed to embed report %d: %s", report_id, e)

    async def create_report(self, session: AsyncSession, report_in: ReportCreate) -> Report:
        r = Report(
            task_id=report_in.task_id,
            title=report_in.title,
            version=report_in.version,
            status=report_in.status,
            summary=report_in.summary,
            content_markdown=report_in.content_markdown,
            storage_ref=report_in.storage_ref,
            images=report_in.chart_images if report_in.chart_images else [],  # NEW: persist chart_images
            pending_vector_upload=True,  # NEW: lazy 模式默认待写入
            # html-ppt Phase 1: 结构化页面 + 主题 + 摘要
            page_model=report_in.page_model if report_in.page_model else None,
            theme=report_in.theme or "minimal-white",
            notes_summary=report_in.notes_summary if report_in.notes_summary else None,
            html_content=report_in.html_content if report_in.html_content else None,
        )
        session.add(r)
        await session.commit()
        await session.refresh(r)
        # 调试日志:验证图片数据是否成功写入 DB
        logger.info(
            "[DEBUG][report] create_report id=%s task_id=%s images_count=%d, page_model=%d, theme=%s",
            r.id, r.task_id, len(r.images or []),
            len(r.page_model or []), r.theme or "minimal-white",
        )
        # 注意：lazy 模式下不在此处触发 Milvus 写入
        return r

    async def upload_to_vector_store_if_pending(self, session: AsyncSession, report_id: int) -> bool:
        """如果报告 pending_vector_upload=True，则嵌入 Milvus 并更新状态。

        触发时机：用户首次导出报告（PDF/DOCX/PPTX/MD 任一格式）或用户上传报告。
        写入成功后将 pending_vector_upload 标记为 False 并记录上传时间。

        Returns:
            True 表示本次确实写入了 Milvus，False 表示已写入过或写入失败。
        """
        from app.services.vectorstore_upload import vs_upload_service
        import datetime

        r = await self.get_report(session, report_id)
        if not r:
            return False
        if not r.pending_vector_upload:
            return False
        chunks = await vs_upload_service.upload_report(r)
        if chunks > 0:
            r.pending_vector_upload = False
            r.vector_uploaded_at = datetime.datetime.utcnow().isoformat()
            await session.commit()
            await session.refresh(r)
            return True
        return False

    async def get_report(self, session: AsyncSession, report_id: int) -> Optional[Report]:
        stmt = select(Report).where(Report.id == report_id)
        res = await session.execute(stmt)
        return res.scalar_one_or_none()

    async def list_reports(
        self, session: AsyncSession, task_id: Optional[str] = None, q: Optional[str] = None, skip: int = 0, limit: int = 20
    ) -> Tuple[List[Report], int]:
        """获取报告列表（分页），并返回符合条件的总记录数。

        返回:
            (报告列表, 总数)
        """
        base_stmt = select(Report)
        if task_id is not None:
            base_stmt = base_stmt.where(Report.task_id == task_id)
        if q:
            base_stmt = base_stmt.where(or_(Report.title.ilike(f"%{q}%"), Report.summary.ilike(f"%{q}%")))

        # 统计总数（应用相同的过滤条件）
        count_stmt = select(func.count()).select_from(base_stmt.subquery())
        total_res = await session.execute(count_stmt)
        total = int(total_res.scalar() or 0)

        # 按创建时间倒序、ID 倒序排序，确保最新报告优先且顺序稳定
        stmt = base_stmt.order_by(Report.created_at.desc(), Report.id.desc()).offset(skip).limit(limit)
        res = await session.execute(stmt)
        return list(res.scalars().all()), total

    async def update_report(self, session: AsyncSession, report_id: int, up: ReportUpdate) -> Optional[Report]:
        r = await self.get_report(session, report_id)
        if not r:
            return None
        data = up.model_dump(exclude_unset=True)
        for k, v in data.items():
            setattr(r, k, v)
        await session.commit()
        await session.refresh(r)
        return r

    async def delete_report(self, session: AsyncSession, report_id: int) -> bool:
        r = await self.get_report(session, report_id)
        if not r:
            return False
        try:
            vs_service = VectorStoreService()
            vs_service.delete_by_report(str(report_id))
        except Exception as e:
            logger.warning("Failed to delete vectors for report %d: %s", report_id, e)
        await session.delete(r)
        await session.commit()
        return True

    async def create_from_workflow(
        self,
        session: AsyncSession,
        task_id: str,
        title: str,
        content_markdown: Optional[str] = None,
        summary: Optional[str] = None,
        chart_images: Optional[List[Dict[str, Any]]] = None,
        dimension_illustrations: Optional[List[Dict[str, Any]]] = None,
        # html-ppt Phase 1: 结构化页面 + 主题 + 摘要
        page_model: Optional[List[Dict[str, Any]]] = None,
        theme: Optional[str] = None,
        notes_summary: Optional[str] = None,
        # html-ppt Phase 2: 完整 HTML 演示文稿
        html_content: Optional[str] = None,
    ) -> Report:
        # 合并 chart_images 和 dimension_illustrations
        all_images = []
        if chart_images:
            all_images.extend(chart_images)
        if dimension_illustrations:
            # dimension_illustrations 格式: [{"section": str, "title": str, "base64": str, "position": int}]
            # 转换为 chart_images 格式并保留 section 和 position 信息
            for ill in dimension_illustrations:
                all_images.append({
                    "title": ill.get("title", ""),
                    "base64": ill.get("base64", ""),
                    "section": ill.get("section", ""),
                    "position": ill.get("position", 0),
                })

        # 调试日志:记录合并后图片数量
        logger.info(
            "[DEBUG][report] create_from_workflow task_id=%s "
            "chart_images=%d, dimension_illustrations=%d, merged_all_images=%d, page_model=%d, theme=%s, html_content=%d",
            task_id,
            len(chart_images or []),
            len(dimension_illustrations or []),
            len(all_images),
            len(page_model or []),
            theme or "minimal-white",
            len(html_content or ""),
        )

        report_in = ReportCreate(
            task_id=task_id,
            title=title,
            status="published",
            content_markdown=content_markdown,
            summary=summary or (content_markdown[:200] if content_markdown else None),
            chart_images=all_images,
            # html-ppt Phase 1
            page_model=page_model,
            theme=theme or "minimal-white",
            notes_summary=notes_summary,
            # html-ppt Phase 2
            html_content=html_content,
        )
        return await self.create_report(session, report_in)

    async def upload_report(
        self,
        session: AsyncSession,
        file_content: bytes,
        filename: str,
        title: Optional[str] = None,
    ) -> Report:
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else "txt"

        if ext == "pdf":
            from app.parsers.pdf import PDFParser
            parser = PDFParser()
            result = await parser.parse(io.BytesIO(file_content), filename)
            content = result.content
        elif ext == "docx":
            from app.parsers.docx import DocxParser
            parser = DocxParser()
            result = await parser.parse(io.BytesIO(file_content), filename)
            content = result.content
        elif ext in ("md", "txt"):
            content = file_content.decode("utf-8", errors="replace")
        else:
            raise ValueError(f"Unsupported file type: .{ext}")

        report_in = ReportCreate(
            title=title or filename,
            status="published",
            content_markdown=content,
            summary=content[:200] if content else None,
        )
        return await self.create_report(session, report_in)

    async def get_statistics(
        self,
        session: AsyncSession,
        period: str,
        report_type: Optional[str] = None,
        status: Optional[str] = None,
        limit: int = 12
    ) -> List[ReportStatisticsItem]:
        """
        获取报告统计数据，按指定时间周期分组统计报告数量

        参数:
            session: 数据库会话
            period: 时间周期 ('day' | 'week' | 'month' | 'year'
            report_type: 预留参数，暂未使用
            status: 预留参数，暂未使用
            limit: 返回最近多少个时间周期

        返回:
            List[ReportStatisticsItem]: 包含标签和数量的统计列表，按时间正序排列
        """
        import datetime
        
        period_map = {
            'day': 'day',
            'week': 'week',
            'month': 'month',
            'year': 'year'
        }
        
        truncated_date = func.date_trunc(period_map[period], Report.created_at).label('period_date')
        
        stmt = (
            select(truncated_date, func.count(Report.id).label('count'))
            .group_by(truncated_date)
            .order_by(truncated_date.desc())
            .limit(limit)
        )
        
        res = await session.execute(stmt)
        results = res.all()
        
        # Build a dict of actual data: label -> count
        actual_data = {}
        for row in results:
            date_obj = row.period_date
            if period == 'day':
                label = date_obj.strftime('%Y-%m-%d')
            elif period == 'week':
                year = date_obj.year
                week = date_obj.isocalendar()[1]
                label = f'{year}-W{week:02d}'
            elif period == 'month':
                label = date_obj.strftime('%Y-%m')
            elif period == 'year':
                label = date_obj.strftime('%Y')
            else:
                label = date_obj.strftime('%Y-%m-%d')
            actual_data[label] = row.count
        
        # Generate complete period labels from now going back 'limit' periods
        now = datetime.datetime.now()
        all_labels = []
        
        for i in range(limit):
            if period == 'day':
                target = now - datetime.timedelta(days=i)
                label = target.strftime('%Y-%m-%d')
            elif period == 'week':
                target = now - datetime.timedelta(weeks=i)
                year = target.year
                week = target.isocalendar()[1]
                label = f'{year}-W{week:02d}'
            elif period == 'month':
                # Calculate month by subtracting i months from now
                month = now.month - i
                year = now.year
                while month <= 0:
                    month += 12
                    year -= 1
                label = f'{year}-{month:02d}'
            elif period == 'year':
                label = str(now.year - i)
            else:
                target = now - datetime.timedelta(days=i)
                label = target.strftime('%Y-%m-%d')
            
            all_labels.append(label)
        
        # Merge: use actual data if available, otherwise 0
        formatted_results = []
        for label in reversed(all_labels):
            count = actual_data.get(label, 0)
            formatted_results.append(ReportStatisticsItem(label=label, count=count))
        
        return formatted_results

    # ------------------------------------------------------------------
    # 配图匹配辅助方法
    # ------------------------------------------------------------------

    @staticmethod
    def _build_section_chart_map(chart_images: List[Dict[str, Any]]) -> Dict[str, List[Dict[str, Any]]]:
        """将 chart_images 按 section 字段分组，返回 {section_name: [chart, ...]}。

        没有 section 字段的配图会被忽略（由调用方走旧逻辑兜底）。
        """
        mapping: Dict[str, List[Dict[str, Any]]] = {}
        if not chart_images:
            return mapping
        for ci in chart_images:
            section = (ci.get("section") or "").strip()
            if not section:
                continue
            mapping.setdefault(section, []).append(ci)
        return mapping

    @staticmethod
    def _match_section(heading: str, section_name: str) -> bool:
        """判断 Markdown 章节标题是否与配图的 section 字段匹配。

        匹配规则：
        1. 完全相等（忽略首尾空白）
        2. 章节标题包含 section 名称（子串匹配，忽略大小写）
        3. section 名称包含章节标题中的关键维度词
        """
        heading_clean = heading.strip()
        section_clean = section_name.strip()
        if not section_clean:
            return False
        if heading_clean == section_clean:
            return True
        if section_clean.lower() in heading_clean.lower():
            return True
        if heading_clean.lower() in section_clean.lower():
            return True
        return False

    @staticmethod
    def _find_charts_for_heading(heading: str, section_chart_map: Dict[str, List[Dict[str, Any]]]) -> List[Dict[str, Any]]:
        """根据章节标题从 section_chart_map 中找到所有匹配的配图。"""
        matched: List[Dict[str, Any]] = []
        for section_name, charts in section_chart_map.items():
            if ReportService._match_section(heading, section_name):
                matched.extend(charts)
        return matched

    @staticmethod
    def _collect_unmatched_charts(
        chart_images: List[Dict[str, Any]],
        section_chart_map: Dict[str, List[Dict[str, Any]]],
        matched_sections: set,
    ) -> List[Dict[str, Any]]:
        """收集未被章节匹配的配图，用于在文档末尾兜底插入。

        包括：
        - 没有 section 字段的配图
        - section 字段存在但未在任何章节中匹配到的配图
        """
        unmatched = []
        if not chart_images:
            return unmatched
        for ci in chart_images:
            section = (ci.get("section") or "").strip()
            if not section:
                unmatched.append(ci)
            elif section not in matched_sections:
                unmatched.append(ci)
        return unmatched

    def _add_chart_to_docx(self, doc, ci: Dict[str, Any], width_inches: float = 5.5):
        """向 docx 文档中插入一张配图。"""
        doc.add_heading(ci.get("title", "图表"), 3)
        try:
            img_data = io.BytesIO(base64.b64decode(ci["base64"]))
            doc.add_picture(img_data, width=docx.shared.Inches(width_inches))
        except Exception as e:
            logger.warning("Failed to insert chart image into docx: %s", e)

    def _add_chart_to_pdf_story(self, story, ci: Dict[str, Any], h2_style, width_px: int = 450):
        """向 PDF story 中插入一张配图,返回创建的临时文件路径列表(用于清理)。"""
        import tempfile
        from PIL import Image as PILImage
        
        story.append(Paragraph(ci.get("title", "图表"), h2_style))
        try:
            img_data = io.BytesIO(base64.b64decode(ci["base64"]))
        except Exception:
            logger.warning("Failed to decode base64 for chart: %s", ci.get("title", "unknown"))
            return []
        
        # 验证图片是否有效
        try:
            img_data.seek(0)
            PILImage.open(img_data).verify()
            img_data.seek(0)
        except Exception as e:
            logger.warning("Invalid image data for chart '%s': %s", ci.get("title", "unknown"), e)
            return []
        
        tmp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
        tmp.write(img_data.getvalue())
        tmp.flush()
        tmp_path = tmp.name
        tmp.close()
        story.append(RLImage(tmp_path, width=width_px, height=int(width_px * 0.6)))
        story.append(Spacer(1, 10))
        return [tmp_path]

    def _add_chart_to_pptx_slide(self, prs, ci: Dict[str, Any]):
        """为单张配图创建一张空白幻灯片并插入图片。"""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        # 添加图表标题
        left = docx.shared.Inches(1)
        top = docx.shared.Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, docx.shared.Inches(8), docx.shared.Inches(0.8))
        txBox.text_frame.text = ci.get("title", "图表")
        txBox.text_frame.paragraphs[0].font.size = docx.shared.Pt(18)
        txBox.text_frame.paragraphs[0].font.bold = True
        # 插入图片
        try:
            img_data = io.BytesIO(base64.b64decode(ci["base64"]))
            slide.shapes.add_picture(
                img_data,
                docx.shared.Inches(1), docx.shared.Inches(1.5),
                width=docx.shared.Inches(8),
            )
        except Exception as e:
            logger.warning("Failed to insert chart image into pptx: %s", e)

    def generate_docx(self, report: Report, chart_images: List[Dict[str, Any]] = None) -> bytes:
        """从 Markdown 报告生成 DOCX 文件，标题层级自动映射为 Word 样式。

        配图插入逻辑：
        - 如果 chart_images 中配图带有 section 字段，则按 section 匹配章节标题，
          在对应章节内容结束后插入配图（宽度 5.5 英寸 ≈ 页面宽度 80%）。
        - 如果配图没有 section 字段或未能匹配到任何章节，则在文档末尾统一插入（向后兼容）。
        """
        doc = docx.Document()
        doc.add_heading(report.title or "未命名报告", 0)
        doc.add_heading("执行摘要", 1)
        doc.add_paragraph(report.summary or "无摘要")
        doc.add_heading("详细内容", 1)

        # 构建 section -> charts 映射
        section_chart_map = self._build_section_chart_map(chart_images)
        matched_sections: set = set()  # 记录已匹配到章节的 section 名称

        # 按章节分割 Markdown 内容
        content = report.content_markdown or ""
        # 使用正则按 ## 级标题分割，保留标题行
        sections = re.split(r"(?=^## )", content, flags=re.MULTILINE)

        for section_text in sections:
            section_text = section_text.strip()
            if not section_text:
                continue

            lines = section_text.split("\n")
            current_heading = None

            for line in lines:
                line_stripped = line.strip()
                if not line_stripped:
                    continue
                if line_stripped.startswith("# "):
                    doc.add_heading(line_stripped[2:], 1)
                elif line_stripped.startswith("## "):
                    heading_text = line_stripped[3:]
                    current_heading = heading_text
                    doc.add_heading(heading_text, 2)
                elif line_stripped.startswith("### "):
                    doc.add_heading(line_stripped[4:], 3)
                else:
                    doc.add_paragraph(line_stripped)

            # 章节内容结束后，检查是否有匹配的配图
            if current_heading and section_chart_map:
                matched_charts = self._find_charts_for_heading(current_heading, section_chart_map)
                if matched_charts:
                    for ci in matched_charts:
                        self._add_chart_to_docx(doc, ci, width_inches=5.5)
                    # 记录已匹配的 section
                    for ci in matched_charts:
                        s = (ci.get("section") or "").strip()
                        if s:
                            matched_sections.add(s)

        # 兜底：未匹配的配图在文档末尾统一插入（向后兼容）
        unmatched_charts = self._collect_unmatched_charts(chart_images, section_chart_map, matched_sections)
        if unmatched_charts:
            doc.add_heading("数据可视化图表", 1)
            for ci in unmatched_charts:
                self._add_chart_to_docx(doc, ci, width_inches=5.5)

        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        return buf.getvalue()

    def generate_pptx(self, report: Report, chart_images: List[Dict[str, Any]] = None, template_id: str = None) -> bytes:
        """从报告生成 PPTX 文件。

        优先使用 PPTTemplateService 基于母版模板填充；当 template_id 为 None 或无
        任何模板可用时回退到旧的"裸版式"生成逻辑，保持向后兼容。

        Args:
            report: 报告 ORM 对象
            chart_images: 配图列表
            template_id: 模板 ID（与 MANIFEST.json 中 id 对应，如 "template1"）
        """
        try:
            svc = PPTTemplateService()
            templates = svc.list_templates()
            chosen_id = template_id
            if not chosen_id and templates:
                chosen_id = templates[0].id
            if chosen_id and svc.get_template(chosen_id):
                return svc.fill_template(chosen_id, report, chart_images)
        except Exception as e:
            logger.warning("PPTTemplateService failed (%s), falling back to legacy pptx", e)

        # 兜底：保留旧版实现（无模板或模板服务失败时使用）
        return self._generate_pptx_legacy(report, chart_images)

    def _generate_pptx_legacy(self, report: Report, chart_images: List[Dict[str, Any]] = None) -> bytes:
        """旧版 PPTX 生成逻辑（无母版模板时使用）

        配图插入逻辑：
        - 对于每个维度章节，创建一张幻灯片包含标题和文本内容
        - 如果该章节有匹配的配图（根据 section 字段），在文本下方插入配图，
          或创建单独的幻灯片展示配图
        - 图片宽度设置为 8 英寸，位置在 (1, 1.5) 英寸处
        - 如果配图没有 section 字段或未能匹配到任何章节，则在文档末尾统一插入（向后兼容）
        """
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = report.title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = report.summary or ""

        # 构建 section -> charts 映射
        section_chart_map = self._build_section_chart_map(chart_images)
        matched_sections: set = set()

        content = report.content_markdown or ""
        sections = re.split(r"(?=^## )", content, flags=re.MULTILINE)
        content_slide_layout = prs.slide_layouts[1]

        for section_text in sections:
            section_text = section_text.strip()
            if not section_text:
                continue

            lines = section_text.split("\n")
            heading = lines[0].strip()
            if heading.startswith("## "):
                heading = heading[3:]
            else:
                # 不是 ## 开头的跳过
                continue

            # 创建内容幻灯片
            slide = prs.slides.add_slide(content_slide_layout)
            slide.shapes.title.text = heading
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()

            # 添加文本内容
            for line in lines[1:]:
                line = line.strip()
                if not line:
                    continue
                p = tf.add_paragraph()
                p.text = line

            # 检查是否有匹配的配图
            if section_chart_map:
                matched_charts = self._find_charts_for_heading(heading, section_chart_map)
                if matched_charts:
                    # 为每张匹配的配图创建单独的幻灯片
                    for ci in matched_charts:
                        self._add_chart_to_pptx_slide(prs, ci)
                    # 记录已匹配的 section
                    for ci in matched_charts:
                        s = (ci.get("section") or "").strip()
                        if s:
                            matched_sections.add(s)

        # 兜底：未匹配的配图在末尾统一插入（向后兼容）
        unmatched_charts = self._collect_unmatched_charts(chart_images, section_chart_map, matched_sections)
        if unmatched_charts:
            for ci in unmatched_charts:
                self._add_chart_to_pptx_slide(prs, ci)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.getvalue()

    def _get_cjk_stylesheet(self):
        """
        获取支持中文的 reportlab 样式表

        为什么需要自定义样式表：
        - reportlab 默认 getSampleStyleSheet() 使用 Helvetica 字体，不含中文字形
        - PDF 绘制中文时发现字体缺少对应 glyph，导致乱码
        - 需要注册中文字体并替换所有样式的 fontName

        字体选择（按优先级探测,任意一个存在即采用）：
        - Linux(Docker): 文泉驿微米黑 (WenQuanYi Micro Hei) — apt-get install fonts-wqy-microhei
        - Windows:      微软雅黑 (msyh.ttc) / 黑体 (simhei.ttf) / 宋体 (simsun.ttc) / 楷体 (simkai.ttf)
        - macOS:        PingFang.ttc / STHeiti Medium.ttc / 苹方
        """
        font_paths = [
            # Linux(Docker 部署环境)
            "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",
            "/usr/share/fonts/wqy-microhei/wqy-microhei.ttc",
            "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc",
            # Windows(本地开发)
            "C:/Windows/Fonts/msyh.ttc",          # 微软雅黑
            "C:/Windows/Fonts/msyhbd.ttc",        # 微软雅黑 Bold
            "C:/Windows/Fonts/simhei.ttf",        # 黑体
            "C:/Windows/Fonts/simsun.ttc",        # 宋体
            "C:/Windows/Fonts/simkai.ttf",        # 楷体
            "C:/Windows/Fonts/STSONG.TTF",        # 华文宋体
            "C:/Windows/Fonts/SIMYOU.TTF",        # 幼圆
            # macOS(本地开发)
            "/System/Library/Fonts/PingFang.ttc",
            "/System/Library/Fonts/STHeiti Medium.ttc",
            "/Library/Fonts/Songti.ttc",
        ]
        cjk_font_name = None
        cjk_font_path = None
        for path in font_paths:
            try:
                if not os.path.exists(path):
                    continue
                pdfmetrics.registerFont(TTFont('WQYMicroHei', path))
                cjk_font_name = 'WQYMicroHei'
                cjk_font_path = path
                break
            except Exception:
                continue

        styles = getSampleStyleSheet()

        if cjk_font_name:
            for style_name in ['Title', 'Heading1', 'Heading2', 'Heading3', 'BodyText']:
                style = styles[style_name]
                style.fontName = cjk_font_name
                style.encoding = 'UTF-8'
        else:
            import logging
            logging.getLogger(__name__).warning(
                "No CJK font found in known paths; PDF will use default font and may show garbled CJK. "
                "Searched: %s",
                font_paths,
            )

        return styles

    def generate_pdf(self, report: Report, chart_images: List[Dict[str, Any]] = None) -> bytes:
        """从 Markdown 报告生成 PDF 文件,支持中文字体。

        配图插入逻辑：
        - 如果 chart_images 中配图带有 section 字段，则按 section 匹配章节标题，
          在对应章节内容结束后插入配图（宽度 450 像素）。
        - 如果配图没有 section 字段或未能匹配到任何章节，则在文档末尾统一插入（向后兼容）。

        鲁棒性说明:
        - Markdown 行经 HTML 转义后再交给 reportlab Paragraph(避免 < / > / & 触发 XMLSyntaxError)
        - 当某行无法被 reportlab 解析时,降级为纯文本 Paragraph
        - 行内的图片/表格语法会被静默跳过(避免垃圾段落)
        """
        from xml.sax.saxutils import escape as _xml_escape

        buf = io.BytesIO()
        doc = SimpleDocTemplate(buf, pagesize=letter, leftMargin=54, rightMargin=54, topMargin=54, bottomMargin=54)
        styles = self._get_cjk_stylesheet()
        t_style = styles['Title']
        h1_style = styles['Heading1']
        h2_style = styles['Heading2']
        body_style = styles['BodyText']

        def _safe_paragraph(text: str, style):
            """把字符串安全地放到 Paragraph;先转义 XML 特殊字符;失败时再二次降级。"""
            if text is None:
                text = ""
            escaped = _xml_escape(text)
            try:
                return Paragraph(escaped, style)
            except Exception:
                return Paragraph(escaped.replace("<", "").replace(">", ""), style)

        story = [_safe_paragraph(report.title or "未命名报告", t_style), Spacer(1, 20)]
        story.append(_safe_paragraph("执行摘要", h1_style))
        story.append(_safe_paragraph(report.summary or "无", body_style))
        story.append(Spacer(1, 15))

        story.append(_safe_paragraph("详细内容", h1_style))

        # 构建 section -> charts 映射
        section_chart_map = self._build_section_chart_map(chart_images)
        matched_sections: set = set()
        tmp_paths: List[str] = []

        # 按章节分割 Markdown 内容
        content = report.content_markdown or ""
        sections = re.split(r"(?=^## )", content, flags=re.MULTILINE)

        for section_text in sections:
            section_text = section_text.strip()
            if not section_text:
                continue

            lines = section_text.split("\n")
            current_heading = None

            for raw_line in lines:
                line = raw_line.strip()
                if not line:
                    continue
                # 跳过 markdown 行内图片语法 ![alt](url)
                if line.startswith("![") and "](http" in line:
                    alt = line.split("]", 1)[0][2:] or "未命名"
                    story.append(_safe_paragraph(f"[图片: {alt}]", body_style))
                    continue
                # 跳过纯水平分割线
                if re.fullmatch(r"[-*_]{3,}", line):
                    continue
                if line.startswith("# "):
                    story.append(Spacer(1, 10))
                    story.append(_safe_paragraph(line[2:], h1_style))
                elif line.startswith("## "):
                    current_heading = line[3:]
                    story.append(Spacer(1, 8))
                    story.append(_safe_paragraph(current_heading, h2_style))
                else:
                    story.append(_safe_paragraph(line, body_style))
                    story.append(Spacer(1, 4))

            # 章节内容结束后，检查是否有匹配的配图
            if current_heading and section_chart_map:
                matched_charts = self._find_charts_for_heading(current_heading, section_chart_map)
                if matched_charts:
                    for ci in matched_charts:
                        paths = self._add_chart_to_pdf_story(story, ci, h2_style, width_px=450)
                        tmp_paths.extend(paths)
                    for ci in matched_charts:
                        s = (ci.get("section") or "").strip()
                        if s:
                            matched_sections.add(s)

        # 兜底：未匹配的配图在文档末尾统一插入（向后兼容）
        unmatched_charts = self._collect_unmatched_charts(chart_images, section_chart_map, matched_sections)
        if unmatched_charts:
            story.append(_safe_paragraph("数据可视化图表", h1_style))
            for ci in unmatched_charts:
                paths = self._add_chart_to_pdf_story(story, ci, h2_style, width_px=450)
                tmp_paths.extend(paths)

        doc.build(story)

        # 清理图表临时文件
        for tp in tmp_paths:
            try:
                os.unlink(tp)
            except OSError:
                pass

        buf.seek(0)
        return buf.getvalue()

    # ------------------------------------------------------------------
    # 专业排版辅助方法（Markdown 排版美化）
    # ------------------------------------------------------------------

    # 章节标题关键词 -> emoji 标记 的映射
    _SECTION_EMOJI_MAP: Dict[str, str] = {
        "执行摘要": "📑",
        "executive summary": "📑",
        "数据": "📊",
        "指标": "📊",
        "统计": "📊",
        "洞察": "💡",
        "建议": "💡",
        "机会": "💡",
        "风险": "⚠️",
        "挑战": "⚠️",
        "威胁": "⚠️",
        "结论": "🎯",
        "总结": "🎯",
        "来源": "🔗",
        "参考": "🔗",
        "竞争格局": "⚔️",
        "趋势": "📈",
        "预测": "🔮",
        "展望": "🔮",
        "细分板块": "🧩",
        "宏观经济": "🌐",
    }

    @classmethod
    def _decorate_heading(cls, heading_text: str) -> str:
        """为章节标题添加 emoji 标记(如果还没有),支持中文/英文标题。

        策略:
        - 已有 emoji(以 unicode 表情符号开头的)则跳过
        - 否则按关键词匹配表查找合适 emoji
        - 兜底: 不修改原标题
        """
        if not heading_text:
            return heading_text
        text = heading_text.strip()
        # 已有 emoji 或特殊符号开头 -> 跳过
        if text and ord(text[0]) >= 0x1F300:
            return heading_text
        # 关键词匹配(不区分大小写)
        text_lower = text.lower()
        for keyword, emoji in cls._SECTION_EMOJI_MAP.items():
            if keyword.lower() in text_lower:
                return f"{emoji} {heading_text}"
        return heading_text

    @classmethod
    def _extract_key_insight_blockquotes(cls, content: str) -> str:
        """把"关键结论"风格的句子(以 【关键】/ **关键** / ⚠️ 等开头)
        包装为 > blockquote 引用块,提升排版专业感。

        启发式规则:
        - 匹配以"**关键**" / "**Key Insight**" / "**Insight**" 开头,到下一个换行或段落结束
        - 整段缩进前缀为 "> ",Markdown 渲染为引用块
        """
        if not content:
            return content
        # 仅作用于未在引用块内的关键句
        pattern = re.compile(
            r'(?m)^(\*\*(?:关键|Key Insight|Insight|关键洞察|核心要点)[^*]*\*\*[^\n]*(?:\n(?!\n)[^\n]*)*)',
        )

        def _to_quote(m: re.Match) -> str:
            block = m.group(1)
            lines = block.split("\n")
            quoted = "\n".join(f"> {line}" for line in lines)
            return quoted

        return pattern.sub(_to_quote, content)

    @classmethod
    def _ensure_section_separators(cls, content: str) -> str:
        """在 ## 章节之间插入水平分隔线 ---,保证排版节奏。

        规则:
        - 跳过 H1 标题前(文档头不加分隔线)
        - 跳过紧邻 --- 的 ## 标题(避免重复)
        """
        if not content:
            return content
        lines = content.split("\n")
        out: List[str] = []
        prev_was_h1 = False
        prev_was_hr = False
        for i, line in enumerate(lines):
            stripped = line.strip()
            if stripped.startswith("## "):
                # 决定是否在前面加 ---
                if not prev_was_h1 and not prev_was_hr and out:
                    # 看 out 末尾是否已有 ---
                    if not (out and out[-1].strip() == "---"):
                        out.append("")
                        out.append("---")
                        out.append("")
            out.append(line)
            prev_was_h1 = stripped.startswith("# ") and not stripped.startswith("## ")
            prev_was_hr = stripped == "---"
        return "\n".join(out)

    @classmethod
    def _enhance_table_formatting(cls, content: str) -> str:
        """改进 Markdown 表格渲染质量。

        功能:
        1. 支持多行表头:连续多行 (---|---) 分隔行(只要|列数一致)会被识别为多行表头
        2. 支持单元格合并标记: `<<span:N>>` 标识该单元格水平合并 N 列
           渲染时展开为对应个数的空 `<td>` (HTML 形式)
           或者在 Markdown 形式下用 `||` 分隔表示合并
        3. 对比表格对齐:自动识别数字列并设置右对齐(原默认|改为|:--:|--:|--:|)
        4. 表格前后自动加空行,保证 Markdown 解析正确

        向后兼容:不包含上述标记的表格保持原样。
        """
        if not content:
            return content

        lines = content.split("\n")
        out: List[str] = []
        i = 0
        while i < len(lines):
            line = lines[i]
            stripped = line.strip()
            # 检测表格起始(以 | 开头且包含 |)
            if stripped.startswith("|") and "|" in stripped[1:] and i + 1 < len(lines):
                next_stripped = lines[i + 1].strip() if i + 1 < len(lines) else ""
                # 表格头分隔行: |---|---|---| 或 |:---|:---:|---:|
                if re.match(r'^\|?\s*:?-{2,}:?\s*(\|\s*:?-{2,}:?\s*)+\|?\s*$', next_stripped):
                    # 抽取整个表格块
                    block = [line, lines[i + 1]]
                    j = i + 2
                    while j < len(lines) and lines[j].strip().startswith("|"):
                        block.append(lines[j])
                        j += 1
                    enhanced = cls._enhance_single_table(block)
                    # 表格前后补空行
                    if out and out[-1].strip() != "":
                        out.append("")
                    out.extend(enhanced)
                    if j < len(lines) and lines[j].strip() != "":
                        out.append("")
                    i = j
                    continue
            out.append(line)
            i += 1
        return "\n".join(out)

    @classmethod
    def _enhance_single_table(cls, block: List[str]) -> List[str]:
        """处理单个 Markdown 表格块(标题行+分隔行+若干数据行)。"""
        if len(block) < 2:
            return block
        header_line = block[0]
        sep_line = block[1]
        data_lines = block[2:]

        # 拆分单元格(支持 | a | b | 和 |a|b| 两种风格)
        def _split_row(row: str) -> List[str]:
            s = row.strip()
            if s.startswith("|"):
                s = s[1:]
            if s.endswith("|"):
                s = s[:-1]
            return [c.strip() for c in s.split("|")]

        def _join_row(cells: List[str]) -> str:
            return "| " + " | ".join(cells) + " |"

        header_cells = _split_row(header_line)
        sep_cells = _split_row(sep_line)
        n_cols = max(len(header_cells), len(sep_cells))

        # === 1) 单元格合并标记 <<span:N>> ===
        # 例如: "指标 <<span:3>>" -> "指标" + 两个空合并格
        def _expand_merges(cells: List[str]) -> List[str]:
            new_cells: List[str] = []
            for c in cells:
                m = re.search(r'<<\s*span\s*:\s*(\d+)\s*>>', c)
                if m:
                    span_n = int(m.group(1))
                    base = re.sub(r'<<\s*span\s*:\s*\d+\s*>>', '', c).strip()
                    new_cells.append(base if base else c)
                    for _ in range(span_n - 1):
                        new_cells.append("↪")  # 用箭头表示"被合并格"
                else:
                    new_cells.append(c)
            return new_cells

        header_cells = _expand_merges(header_cells)
        data_lines = [_expand_merges(_split_row(dl)) for dl in data_lines]

        # === 2) 对比表格自动对齐 ===
        # 判断规则:
        # - 如果数据行中超过 50% 的单元格是"数字 + 可选单位",则该列右对齐
        # - 如果表头含"占比/比例/率/%/增长",则该列右对齐
        def _is_numeric_cell(s: str) -> bool:
            t = s.strip().replace(",", "").replace(" ", "")
            # 数字 + 可选 %, ‰, 元, 亿, 万, $, ¥
            return bool(re.match(r'^-?\d+(?:\.\d+)?\s*[%‰元亿元万元$/¥]?$', t))

        align_keywords = ("占比", "比例", "率", "%", "增长", "增速", "数量", "金额", "规模", "rank", "排名")
        aligns: List[str] = []
        for col_idx, h in enumerate(header_cells):
            col_data = [dl[col_idx] if col_idx < len(dl) else "" for dl in data_lines]
            if any(kw in h for kw in align_keywords):
                aligns.append("right")
            elif col_data and sum(1 for v in col_data if _is_numeric_cell(v)) / max(len(col_data), 1) >= 0.5:
                aligns.append("right")
            else:
                aligns.append("left")

        def _align_cell(align: str) -> str:
            if align == "right":
                return "---:"
            return "---"

        # 重新生成分隔行
        # 处理: 列数不一致时(合并后)补充 --- 或删减
        if len(sep_cells) < n_cols:
            sep_cells = sep_cells + ["---"] * (n_cols - len(sep_cells))
        elif len(sep_cells) > n_cols:
            sep_cells = sep_cells[:n_cols]
        new_sep_cells = [_align_cell(a) for a in aligns[:n_cols]]
        if len(new_sep_cells) < n_cols:
            new_sep_cells = new_sep_cells + ["---"] * (n_cols - len(new_sep_cells))
        new_sep_line = _join_row(new_sep_cells)

        # 同样补齐 header / data 行
        def _pad(cells: List[str]) -> List[str]:
            if len(cells) < n_cols:
                return cells + [""] * (n_cols - len(cells))
            return cells[:n_cols]

        new_header = _join_row(_pad(header_cells))
        new_data = [_join_row(_pad(dl)) for dl in data_lines]

        # === 3) 多行表头支持 ===
        # 如果 header_line 包含多行(以 \n 分隔,在同一字符串内),分别处理
        # 简单场景: 用户把多行 header 写在一行内并用 <br> 分隔
        if "<br>" in new_header or "<br/>" in new_header:
            br_pat = re.compile(r'<br\s*/?>')
            br_lines = br_pat.split(new_header)
            br_cells_list = [_split_row(l) for l in br_lines if l.strip()]
            # 重新组装为多行 header
            multi_header = "\n".join(_join_row(_pad(c)) for c in br_cells_list)
            return [multi_header, new_sep_line] + new_data

        return [new_header, new_sep_line] + new_data

    def _enhance_markdown(self, content: str) -> str:
        """综合排版美化:
        1. 关键句 -> blockquote
        2. ## 章节之间插入 --- 分隔线
        3. 表格增强(多行表头 / 单元格合并 / 对齐)
        4. 章节标题添加 emoji 标记
        """
        if not content:
            return content
        # 表格增强先做(避免后续插入 \n 干扰表格解析)
        content = self._enhance_table_formatting(content)
        # blockquote 包装
        content = self._extract_key_insight_blockquotes(content)
        # 章节分隔线
        content = self._ensure_section_separators(content)
        # 标题 emoji 装饰
        def _decorate_h(m: re.Match) -> str:
            hashes = m.group(1)
            text = m.group(2)
            return f"{hashes} {self._decorate_heading(text)}"
        content = re.sub(r'^(#{1,6})\s+(.+)$', _decorate_h, content, flags=re.MULTILINE)
        return content

    def generate_markdown(self, report: Report) -> bytes:
        """生成 Markdown 报告(供下载或直接展示)。

        与旧版本的差异:
        - 加入排版美化: 关键句包装为 blockquote、章节之间加分隔线、表格增强
        - 章节标题自动添加 emoji 标记(📊 数据, 💡 洞察, ⚠️ 风险)
        - 向后兼容: 当 content_markdown 为空时,降级为旧版简单拼接
        """
        summary = report.summary or "无"
        content = report.content_markdown or ""
        try:
            enhanced_content = self._enhance_markdown(content) if content else ""
        except Exception as e:
            logger.warning("Markdown enhancement failed, falling back to raw content: %s", e)
            enhanced_content = content
        body = f"# {report.title}\n\n## 📑 执行摘要\n{summary}\n\n---\n\n## 详细内容\n{enhanced_content or content}"
        return body.encode("utf-8")

    async def export_report_with_html_pipeline(
        self, session: AsyncSession, report_id: int, fmt: str, 
        page_model: ReportPageModel, template_id: Optional[str] = None
    ) -> Tuple[str, str, bytes]:
        """
        使用 HTML 生成流程导出报告
        
        新流程：
        1. 将 ReportPageModel 转换为 HTML
        2. 根据目标格式，将 HTML 转换为对应格式
        
        Args:
            session: 数据库会话
            report_id: 报告 ID
            fmt: 导出格式（pptx/pdf/docx/md）
            page_model: 带有布局信息的 ReportPageModel
            template_id: PPT 模板 ID（可选）
            
        Returns:
            (文件名, MIME类型, 文件二进制内容)
        """
        r = await self.get_report(session, report_id)
        if not r:
            raise ValueError("Report not found")
        
        fmt = fmt.lower()
        
        # 1. 将 ReportPageModel 转换为 HTMLPageModel 并生成 HTML
        from app.services.html_report_generator import (
            HTMLReportGenerator, HTMLPageModel, HTMLTextBlock, HTMLImageBlock, LayoutType,
        )
        html_pages = self._convert_pages_to_html(page_model.pages)
        generator = HTMLReportGenerator()
        html_content = generator.generate(html_pages)
        
        logger.info(f"HTML generated for report {report_id}: {len(html_content)} chars")
        
        # 2. 根据格式转换为对应格式
        if fmt == "pptx":
            from app.services.html_to_ppt_converter import HTMLToPPTConverter
            converter = HTMLToPPTConverter()
            
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
                tmp_path = tmp.name
            
            try:
                await converter.convert(html_content, tmp_path, template_id)
                with open(tmp_path, 'rb') as f:
                    data = f.read()
            finally:
                if os.path.exists(tmp_path):
                    os.unlink(tmp_path)
            
            result = (
                f"report_{r.id}.pptx",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                data
            )
        
        elif fmt == "pdf":
            from app.services.html_to_pdf_converter import HTMLToPDFConverter
            converter = HTMLToPDFConverter()
            
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
                tmp_path = tmp.name
            
            try:
                await converter.convert(html_content, tmp_path)
                with open(tmp_path, 'rb') as f:
                    data = f.read()
            finally:
                if os.path.exists(tmp_path):
                    os.unlink(tmp_path)
            
            result = (
                f"report_{r.id}.pdf",
                "application/pdf",
                data
            )
        
        elif fmt == "docx":
            from app.services.html_to_docx_converter import HTMLToDOCXConverter
            converter = HTMLToDOCXConverter()
            
            with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
                tmp_path = tmp.name
            
            try:
                await converter.convert(html_content, tmp_path)
                with open(tmp_path, 'rb') as f:
                    data = f.read()
            finally:
                if os.path.exists(tmp_path):
                    os.unlink(tmp_path)
            
            result = (
                f"report_{r.id}.docx",
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                data
            )
        
        elif fmt in ("md", "markdown"):
            from app.services.html_to_markdown_converter import HTMLToMarkdownConverter
            converter = HTMLToMarkdownConverter()
            
            with tempfile.NamedTemporaryFile(suffix='.md', delete=False, mode='w', encoding='utf-8') as tmp:
                tmp_path = tmp.name
            
            try:
                await converter.convert(html_content, tmp_path)
                with open(tmp_path, 'r', encoding='utf-8') as f:
                    data = f.read().encode('utf-8')
            finally:
                if os.path.exists(tmp_path):
                    os.unlink(tmp_path)
            
            result = (
                f"report_{r.id}.md",
                "text/markdown",
                data
            )
        
        else:
            raise ValueError(f"Unsupported format: {fmt}")
        
        logger.info(f"Report {report_id} exported as {fmt} using HTML pipeline")
        return result
    
    def _choose_layout(self, page: PageModel) -> "LayoutType":
        """单点布局决策 —— 整条 HTML 流程中唯一决定 layout 的地方

        决策优先级（从高到低）：
        1. page_type 显式标记 (cover / section / toc)
        2. 页面内容特征（tables / images 数量 / text_blocks 数量）
        3. factory 默认 layout_hint 兜底

        LayoutType 与 html-ppt 布局的对应关系：
        - cover / section  →  COVER (渐变大标题) / SECTION (章节封面)
        - toc              →  TOC (目录卡片网格)
        - kpi 指标         →  KPI_GRID
        - 单张大图         →  IMAGE_HERO
        - 多图             →  IMAGE_GRID
        - 表格             →  TABLE
        - 图文             →  CONTENT (左文右图)
        """
        from app.services.html_report_generator import LayoutType

        # 1) page_type 优先 —— 业务语义强
        if page.page_type == "cover":
            return LayoutType.COVER
        if page.page_type == "section":
            return LayoutType.SECTION
        if page.page_type == "toc":
            return LayoutType.TOC
        if page.page_type == "summary":
            return LayoutType.THREE_COLUMN  # 总结页用三栏

        # 2) 内容特征
        if page.tables:
            return LayoutType.TABLE
        # 检测文本块里内嵌的 markdown 表格（"| --- |" 形式）
        has_inline_table = any(
            "|" in (tb.text or "") and "---" in (tb.text or "")
            for tb in page.text_blocks
        )
        if has_inline_table:
            return LayoutType.TABLE

        image_count = len(page.images)
        text_block_count = len(page.text_blocks)
        if image_count >= 4:
            return LayoutType.IMAGE_GRID
        if image_count >= 1 and text_block_count <= 1:
            return LayoutType.IMAGE_HERO
        if image_count > 0 and text_block_count <= 2:
            return LayoutType.TWO_COLUMN
        if image_count > 0 and text_block_count > 2:
            return LayoutType.CONTENT

        # 3) factory 默认 layout_hint 兜底
        hint_fallback = {
            "image_only": LayoutType.IMAGE_GRID,
            "text_left": LayoutType.TWO_COLUMN,
            "text_only": LayoutType.CONTENT,
            "text_top": LayoutType.CONTENT,
        }
        return hint_fallback.get(
            page.layout_hint, LayoutType.CONTENT,
        )

    def _convert_pages_to_html(self, pages: List[PageModel]) -> List:
        """将 report_page_model.PageModel 列表转换为 html_report_generator.HTMLPageModel 列表

        Args:
            pages: report_page_model.PageModel 列表

        Returns:
            html_report_generator.HTMLPageModel 列表
        """
        from app.services.html_report_generator import (
            HTMLPageModel, HTMLTextBlock, HTMLImageBlock, LayoutType,
        )

        html_pages = []
        for page in pages:
            # 转换文本块（过滤掉与 page.title 重复的 TextBlock，
            # 因为 make_content_page / make_summary_page 会把标题写入 text_blocks[0]，
            # 而 HTML 渲染器已经用 <h2>{page.title}</h2> 单独渲染标题）
            text_blocks = []
            for tb in page.text_blocks:
                if tb.text.strip() == page.title.strip():
                    continue
                text_blocks.append(HTMLTextBlock(
                    text=tb.text,
                    emphasis=[],  # PageModel 没有 emphasis 字段，留空
                ))

            # 转换图片块
            image_blocks = []
            for img in page.images:
                # PageModel.ImageBlock 使用 base64，HTMLImageBlock 使用 url
                # 将 base64 转换为 data URL
                if img.base64:
                    url = f"data:image/png;base64,{img.base64}"
                    image_blocks.append(HTMLImageBlock(
                        url=url,
                        caption=img.alt or "",
                        source="matplotlib",
                    ))

            # 转换表格数据（取第一个表格；HTML 渲染器目前只支持单表格）
            table_data = None
            if page.tables:
                tbl = page.tables[0]
                table_data = {
                    "headers": tbl.headers,
                    "rows": tbl.rows,
                }

            # 单点布局决策（不再依赖 _enhance_layout_info 写入的 layout_hint 字符串）
            layout = self._choose_layout(page)

            # 生成 kicker：封面用 subtitle，section 页用编号，其他留空
            kicker = ""
            if page.page_type == "cover":
                kicker = page.subtitle
            elif page.page_type == "section":
                kicker = f"{page.section_number:02d} / {page.title}"

            html_page = HTMLPageModel(
                title=page.title,
                layout=layout,
                kicker=kicker,
                text_blocks=text_blocks,
                image_blocks=image_blocks,
                kpi_metrics=[],
                table_data=table_data,
            )
            html_pages.append(html_page)

        return html_pages
    
    async def export_report(
        self, session: AsyncSession, report_id: int, fmt: str, template_id: Optional[str] = None
    ) -> Tuple[str, str, bytes]:
        """
        统一导出入口

        返回三元组：(文件名, MIME类型, 文件二进制内容)
        路由函数据此设置 HTTP 响应头（Content-Disposition + Content-Type）

        Args:
            template_id: PPT 母版模板 ID（仅 fmt=pptx 时生效），None 时取第一套可用模板

        注意：lazy 模式下不在此处触发 Milvus 写入；
        写入由路由层在调用本方法前通过 `upload_to_vector_store_if_pending()` 触发。
        """
        r = await self.get_report(session, report_id)
        if not r:
            raise ValueError("Report not found")

        chart_images = r.images if isinstance(r.images, list) else []
        # 调试日志:验证从 DB 读出的图片数量
        logger.info(
            "[DEBUG][report] export_report id=%s fmt=%s images_from_db=%d",
            report_id, fmt, len(chart_images),
        )

        fmt = fmt.lower()
        if fmt == "docx":
            data = self.generate_docx(r, chart_images=chart_images)
            result = (f"report_{r.id}.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document", data)
        elif fmt == "pdf":
            data = self.generate_pdf(r, chart_images=chart_images)
            result = (f"report_{r.id}.pdf", "application/pdf", data)
        elif fmt == "pptx":
            data = self.generate_pptx(r, chart_images=chart_images, template_id=template_id)
            result = (f"report_{r.id}.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation", data)
        elif fmt == "md" or fmt == "markdown":
            result = (f"report_{r.id}.md", "text/markdown", self.generate_markdown(r))
        else:
            raise ValueError("Unsupported format. Use docx, pdf, pptx, or md.")

        return result

    # ── 新接口：基于 ReportPageModel 的统一导出 ───────────────

    def generate_docx_from_model(self, model: ReportPageModel) -> bytes:
        """从 ReportPageModel 生成 DOCX，每页对应一个 Word 页面"""
        doc = docx.Document()

        # 封面
        cover_pages = model.get_pages_by_type("cover")
        if cover_pages:
            cover = cover_pages[0]
            doc.add_heading(model.title, 0)
            subtitle = cover.subtitle or ""
            if subtitle:
                doc.add_paragraph(subtitle)
            doc.add_page_break()

        # 逐页渲染
        for page in model.pages:
            if page.page_type == "cover":
                continue  # 封面已处理

            # 标题
            title_block = next((tb for tb in page.text_blocks if tb.style == "title"), None)
            if title_block:
                doc.add_heading(title_block.text, level=2)

            # 正文
            body_blocks = [tb for tb in page.text_blocks if tb.style != "title"]
            for tb in body_blocks:
                if tb.style == "caption":
                    p = doc.add_paragraph()
                    p.style = doc.styles['Normal']
                    run = p.add_run(tb.text)
                    run.font.size = docx.shared.Pt(tb.font_size)
                    run.font.italic = True
                elif tb.style == "bullet":
                    p = doc.add_paragraph(tb.text, style='List Bullet')
                else:
                    p = doc.add_paragraph(tb.text)
                    p.style = doc.styles['Normal']

            # 图片
            for img in page.images:
                try:
                    img_data = io.BytesIO(base64.b64decode(img.base64.strip()))
                    doc.add_picture(img_data, width=docx.shared.Inches(5.5))
                except Exception as e:
                    logger.warning("Failed to insert image in DOCX: %s", e)

            doc.add_page_break()

        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        return buf.getvalue()

    def generate_pdf_from_model(self, model: ReportPageModel) -> bytes:
        """从 ReportPageModel 生成 PDF，每页对应一个 PDF 页面"""
        from xml.sax.saxutils import escape as _xml_escape

        buf = io.BytesIO()
        doc = SimpleDocTemplate(buf, pagesize=letter,
                                leftMargin=54, rightMargin=54,
                                topMargin=54, bottomMargin=54)
        styles = self._get_cjk_stylesheet()
        t_style = styles['Title']
        h2_style = styles['Heading2']
        body_style = styles['BodyText']

        def _safe_para(text, style):
            if not text:
                text = ""
            escaped = _xml_escape(text)
            try:
                return Paragraph(escaped, style)
            except Exception:
                return Paragraph(escaped.replace("<", "").replace(">", ""), style)

        story = []
        tmp_paths: List[str] = []

        # 封面
        cover_pages = model.get_pages_by_type("cover")
        if cover_pages:
            cover = cover_pages[0]
            story.append(_safe_para(model.title, t_style))
            story.append(Spacer(1, 20))
            if cover.subtitle:
                story.append(_safe_para(cover.subtitle, body_style))
                story.append(Spacer(1, 15))

        # 逐页渲染
        for page in model.pages:
            if page.page_type == "cover":
                continue

            story.append(Spacer(1, 10))

            # 标题
            title_block = next((tb for tb in page.text_blocks if tb.style == "title"), None)
            if title_block:
                story.append(_safe_para(title_block.text, h2_style))
                story.append(Spacer(1, 4))

            # 正文
            body_blocks = [tb for tb in page.text_blocks if tb.style != "title"]
            for tb in body_blocks:
                if tb.style == "caption":
                    story.append(_safe_para(f"[{tb.text}]", body_style))
                else:
                    story.append(_safe_para(tb.text, body_style))
                story.append(Spacer(1, 4))

            # 图片
            for img in page.images:
                try:
                    img_data = io.BytesIO(base64.b64decode(img.base64.strip()))
                    from PIL import Image as PILImage
                    img_data.seek(0)
                    PILImage.open(img_data).verify()
                    img_data.seek(0)

                    tmp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
                    tmp.write(img_data.getvalue())
                    tmp.flush()
                    tmp_path = tmp.name
                    tmp.close()
                    story.append(RLImage(tmp_path, width=450, height=int(450 * 0.6)))
                    story.append(Spacer(1, 10))
                    tmp_paths.append(tmp_path)
                except Exception as e:
                    logger.warning("Failed to insert image in PDF: %s", e)

            story.append(PageBreak())

        # 移除最后一个多余的 PageBreak
        if story and isinstance(story[-1], PageBreak):
            story.pop()

        doc.build(story)

        for tp in tmp_paths:
            try:
                os.unlink(tp)
            except OSError:
                pass

        buf.seek(0)
        return buf.getvalue()

    def generate_pptx_from_model(
        self, model: ReportPageModel, template_id: str = None
    ) -> bytes:
        """从 ReportPageModel 生成 PPTX"""
        try:
            svc = PPTTemplateService()
            templates = svc.list_templates()
            chosen_id = template_id
            if not chosen_id and templates:
                chosen_id = templates[0].id
            if chosen_id and svc.get_template(chosen_id):
                return svc.fill_template_from_model(chosen_id, model)
        except Exception as e:
            logger.warning("PPTTemplateService failed (%s), falling back", e)

        return self._generate_pptx_from_model_fallback(model)

    def _generate_pptx_from_model_fallback(self, model: ReportPageModel) -> bytes:
        """无模板 PPTX 兜底"""
        from pptx import Presentation as DefaultPresentation
        from pptx.util import Pt as PptxPt
        prs = DefaultPresentation()
        blank_layout = prs.slide_layouts[6]

        for page in model.pages:
            slide = prs.slides.add_slide(blank_layout)
            title_block = next((tb for tb in page.text_blocks if tb.style == "title"), None)
            if title_block:
                tx = slide.shapes.add_textbox(
                    docx.shared.Inches(0.5), docx.shared.Inches(0.3),
                    docx.shared.Inches(9.0), docx.shared.Inches(0.8),
                )
                tf = tx.text_frame
                p = tf.paragraphs[0]
                p.text = title_block.text
                p.font.size = PptxPt(24)
                p.font.bold = True

            body_blocks = [tb for tb in page.text_blocks if tb.style != "title"]
            if body_blocks:
                tx = slide.shapes.add_textbox(
                    docx.shared.Inches(0.5), docx.shared.Inches(1.5),
                    docx.shared.Inches(9.0), docx.shared.Inches(5.0),
                )
                tf = tx.text_frame
                tf.word_wrap = True
                for i, tb in enumerate(body_blocks):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = tb.text
                    p.font.size = PptxPt(tb.font_size)

            if page.images:
                try:
                    img = page.images[0]
                    img_blob = io.BytesIO(base64.b64decode(img.base64.strip()))
                    slide.shapes.add_picture(
                        img_blob,
                        docx.shared.Inches(1), docx.shared.Inches(4.5),
                        width=docx.shared.Inches(8),
                    )
                except Exception:
                    pass

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.getvalue()

    async def export_report_from_model(
        self,
        model: ReportPageModel,
        fmt: str,
        template_id: Optional[str] = None,
        report_id: int = 0,
    ) -> Tuple[str, str, bytes]:
        """基于 ReportPageModel 的统一导出入口

        Returns:
            (文件名, MIME类型, 文件二进制内容)
        """
        fmt = fmt.lower()
        if fmt == "docx":
            data = self.generate_docx_from_model(model)
            return (f"report_{report_id}.docx",
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    data)
        elif fmt == "pdf":
            data = self.generate_pdf_from_model(model)
            return (f"report_{report_id}.pdf", "application/pdf", data)
        elif fmt == "pptx":
            data = self.generate_pptx_from_model(model, template_id=template_id)
            return (f"report_{report_id}.pptx",
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    data)
        elif fmt in ("md", "markdown"):
            # Markdown 直接从 model 重建
            md = self._generate_markdown_from_model(model)
            return (f"report_{report_id}.md", "text/markdown", md.encode("utf-8"))
        else:
            raise ValueError("Unsupported format. Use docx, pdf, pptx, or md.")

    def _generate_markdown_from_model(self, model: ReportPageModel) -> str:
        """从 ReportPageModel 重建 Markdown"""
        lines = []
        for page in model.pages:
            title_block = next((tb for tb in page.text_blocks if tb.style == "title"), None)
            if title_block:
                if page.page_type == "cover":
                    lines.append(f"# {title_block.text}")
                elif page.page_type == "section":
                    lines.append(f"## {title_block.text}")
                elif page.page_type == "content":
                    lines.append(f"### {title_block.text}")
                else:
                    lines.append(f"## {title_block.text}")
                lines.append("")

            body_blocks = [tb for tb in page.text_blocks if tb.style != "title"]
            for tb in body_blocks:
                if tb.style == "bullet":
                    lines.append(f"- {tb.text}")
                elif tb.style == "caption":
                    lines.append(f"> {tb.text}")
                else:
                    lines.append(tb.text)

            for img in page.images:
                lines.append(f"![{img.alt}](data:image/png;base64,{img.base64[:50]}...)")

            lines.append("")
            lines.append("---")
            lines.append("")

        return "\n".join(lines)

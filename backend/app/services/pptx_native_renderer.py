"""
PPTX Native Renderer —— 原生 PPTX 幻灯片渲染器

对标豆包/千问/Gamma 的方案：用 python-pptx 创建可编辑的原生 PPT 元素，
而不是截图后插入位图。

与旧版 html_to_ppt_converter.py 的区别：
  旧版：HTML → Playwright 截图 → 插入 PPT（位图，不可编辑）
  新版：HTMLPageModel → python-pptx 原生元素（可编辑文本、可移动图片）

支持的 PPT 模板：
  如果指定 template_path，会从模板 PPT 读取母版（slide master），
  渲染时使用模板的配色方案和字体设置。

布局策略：
  - 封面页：居中标题 + 副标题
  - 目录页：编号卡片网格
  - 章节页：全屏渐变大标题
  - 内容页：根据 LayoutEngine 决定的 variant 选择左右/上下/双栏
  - KPI 页：大数字卡片
  - 图表页：原生 Chart 对象（可编辑数据）
  - 图片页：原生 Picture Shape
  - 表格页：原生 Table
"""

import base64
import io
import logging
import os
import tempfile
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

from app.services.html_report_generator import (
    HTMLPageModel, HTMLTextBlock, HTMLImageBlock, LayoutType,
)
from app.services.layout_engine import (
    LayoutEngine, LayoutDecision, TemplateVariant,
    layout_engine as default_layout_engine,
)

logger = logging.getLogger(__name__)

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN = Inches(0.6)
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN * 2
CONTENT_HEIGHT = SLIDE_HEIGHT - Inches(1.2)  # 扣除 footer

# Color palette (consistent with html-ppt themes)
COLORS = {
    "text_primary": RGBColor(0x11, 0x12, 0x16),
    "text_secondary": RGBColor(0x55, 0x59, 0x6A),
    "text_tertiary": RGBColor(0x8A, 0x8F, 0x9E),
    "accent": RGBColor(0x3B, 0x6C, 0xFF),
    "accent_2": RGBColor(0x7A, 0x5C, 0xFF),
    "bg": RGBColor(0xFF, 0xFF, 0xFF),
    "bg_soft": RGBColor(0xF5, 0xF5, 0xF6),
    "border": RGBColor(0xE0, 0xE0, 0xE5),
    "good": RGBColor(0x1A, 0xAF, 0x6C),
    "bad": RGBColor(0xE0, 0x44, 0x5A),
    "warn": RGBColor(0xF5, 0xA5, 0x24),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
}

# Font settings
FONT_FAMILY = "Microsoft YaHei"  # 微软雅黑（Windows 上最可靠的中文字体）
FONT_FAMILY_EN = "Arial"


class PPTXNativeRenderer:
    """原生 PPTX 幻灯片渲染器"""

    def __init__(
        self,
        layout_engine: Optional[LayoutEngine] = None,
        template_path: Optional[str] = None,
    ):
        self.layout_engine = layout_engine or default_layout_engine
        self.template_path = template_path

    def render(
        self,
        pages: List[HTMLPageModel],
        output_path: str,
        title: str = "Market Insight Report",
    ) -> str:
        """将页面列表渲染为原生 PPTX 文件

        Args:
            pages: HTML 页面模型列表
            output_path: PPTX 输出路径
            title: 演示文稿标题

        Returns:
            PPTX 文件路径
        """
        # 创建 presentation（优先使用模板）
        if self.template_path and os.path.exists(self.template_path):
            prs = Presentation(self.template_path)
            logger.info(f"Using template: {self.template_path}")
        else:
            prs = Presentation()
            prs.slide_width = SLIDE_WIDTH
            prs.slide_height = SLIDE_HEIGHT
            logger.info("Using default blank presentation (no template)")

        total = len(pages)
        for idx, page in enumerate(pages):
            try:
                # LayoutEngine 决策
                decision = self.layout_engine.decide(page)

                # 渲染单页
                self._render_page(prs, page, decision, idx, total)

                # 渲染溢出页
                for overflow_page in decision.overflow_pages:
                    overflow_decision = self.layout_engine.decide(overflow_page)
                    self._render_page(prs, overflow_page, overflow_decision, idx, total)
            except Exception as e:
                logger.warning(f"Failed to render slide {idx}: {e}, using fallback")
                self._render_fallback(prs, page, idx, total)

        # 保存
        prs.save(output_path)
        logger.info(f"Native PPTX saved: {output_path} ({len(prs.slides)} slides)")
        return output_path

    def _render_page(
        self,
        prs: Presentation,
        page: HTMLPageModel,
        decision: LayoutDecision,
        idx: int,
        total: int,
    ):
        """根据布局决策渲染单页幻灯片"""
        # 添加空白幻灯片
        blank_layout = prs.slide_layouts[6]  # index 6 = blank layout
        slide = prs.slides.add_slide(blank_layout)

        # 设置背景
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS["bg"]

        # 按布局类型分发
        renderers = {
            LayoutType.COVER: self._render_cover_page,
            LayoutType.TOC: self._render_toc_page,
            LayoutType.SECTION: self._render_section_page,
            LayoutType.CONTENT: self._render_content_page,
            LayoutType.BULLETS: self._render_bullets_page,
            LayoutType.KPI_GRID: self._render_kpi_grid_page,
            LayoutType.TWO_COLUMN: self._render_two_column_page,
            LayoutType.THREE_COLUMN: self._render_three_column_page,
            LayoutType.TABLE: self._render_table_page,
            LayoutType.IMAGE_HERO: self._render_image_hero_page,
            LayoutType.IMAGE_GRID: self._render_image_grid_page,
            LayoutType.STAT_HIGHLIGHT: self._render_stat_page,
            LayoutType.THANKS: self._render_thanks_page,
            LayoutType.CHART_BAR: self._render_chart_page,
            LayoutType.CHART_LINE: self._render_chart_page,
            LayoutType.CHART_PIE: self._render_chart_page,
        }
        renderer = renderers.get(page.layout, self._render_content_page)
        renderer(slide, page, decision)

        # 添加页脚
        self._add_footer(slide, page, idx, total)

    # ══════════════════════════════════════════════
    # 封面页
    # ══════════════════════════════════════════════
    def _render_cover_page(self, slide, page: HTMLPageModel, decision: LayoutDecision):
        """封面页：居中标题 + 副标题 + 标签"""
        center_x = SLIDE_WIDTH // 2
        title_top = Inches(2.0)
        title_width = Inches(10)
        title_height = Inches(1.8)

        # 主标题
        tf = self._add_text_box(
            slide, center_x - title_width // 2, title_top,
            title_width, title_height,
        )
        p = tf.paragraphs[0]
        self._set_run(p, page.title, size=44, bold=True, color=COLORS["text_primary"], alignment=PP_ALIGN.CENTER)

        # Kicker / 副标题
        if page.kicker:
            tf2 = self._add_text_box(
                slide, center_x - title_width // 2, Inches(0.8),
                title_width, Inches(0.6),
            )
            p2 = tf2.paragraphs[0]
            self._set_run(p2, page.kicker, size=14, color=COLORS["accent"], alignment=PP_ALIGN.CENTER)

        # lede 简介
        if page.text_blocks:
            lede_text = page.text_blocks[0].text[:200]
            tf3 = self._add_text_box(
                slide, center_x - Inches(5), Inches(4.2),
                Inches(10), Inches(1.5),
            )
            p3 = tf3.paragraphs[0]
            self._set_run(p3, lede_text, size=16, color=COLORS["text_secondary"], alignment=PP_ALIGN.CENTER)

        # Pills / 标签（text_blocks[1:4]）
        pill_top = Inches(5.5)
        for i, tb in enumerate(page.text_blocks[1:4]):
            pill_left = center_x - Inches(2.5) + Inches(i * 2.0)
            pill = self._add_rounded_rect(
                slide, pill_left, pill_top, Inches(1.6), Inches(0.45),
                fill_color=COLORS["accent"],
            )
            self._set_run(pill.paragraphs[0], tb.text[:20], size=11, color=COLORS["white"], alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════
    # 目录页
    # ══════════════════════════════════════════════
    def _render_toc_page(self, slide, page: HTMLPageModel, decision: LayoutDecision):
        """目录页：编号卡片网格"""
        items = page.text_blocks[:12]
        cols = 2
        card_height = Inches(1.0)
        card_width = Inches(5.5)
        start_left = MARGIN
        start_top = Inches(2.0)
        gap_x = Inches(0.3)
        gap_y = Inches(0.2)

        # 标题
        tf = self._add_text_box(slide, MARGIN, Inches(0.5), CONTENT_WIDTH, Inches(0.8))
        p = tf.paragraphs[0]
        self._set_run(p, page.title or "目录", size=36, bold=True, color=COLORS["text_primary"])

        for i, tb in enumerate(items):
            col = i % cols
            row = i // cols
            left = start_left + col * (card_width + gap_x)
            top = start_top + row * (card_height + gap_y)

            # 卡片背景
            card = slide.shapes.add_shape(
                1,  # MSO_SHAPE.RECTANGLE
                left, top, card_width, card_height,
            )
            card.fill.solid()
            card.fill.fore_color.rgb = COLORS["bg_soft"]
            card.line.fill.background()

            # 编号
            num_tf = card.text_frame
            num_tf.word_wrap = True
            p_num = num_tf.paragraphs[0]
            self._set_run(p_num, f"{i + 1:02d}", size=20, bold=True, color=COLORS["accent"])

            # 标题
            p_title = num_tf.add_paragraph()
            self._set_run(p_title, tb.text[:40], size=16, color=COLORS["text_primary"])

    # ══════════════════════════════════════════════
    # 章节分隔页
    # ══════════════════════════════════════════════
    def _render_section_page(self, slide, page: HTMLPageModel, decision: LayoutDecision):
        """章节分隔页：居中大标题"""
        center_x = SLIDE_WIDTH // 2
        center_y = SLIDE_HEIGHT // 2

        # 大标题
        tf = self._add_text_box(
            slide, center_x - Inches(5), center_y - Inches(1.2),
            Inches(10), Inches(2.4),
        )
        p = tf.paragraphs[0]
        self._set_run(p, page.title, size=54, bold=True, color=COLORS["accent"], alignment=PP_ALIGN.CENTER)

        # 分隔线
        line = slide.shapes.add_shape(
            1, center_x - Inches(2), center_y + Inches(1.5),
            Inches(4), Pt(4),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = COLORS["accent"]
        line.line.fill.background()

    # ══════════════════════════════════════════════
    # 内容页（核心：支持图文混排）
    # ══════════════════════════════════════════════
    def _render_content_page(self, slide, page: HTMLPageModel, decision: LayoutDecision):
        """智能内容页：根据 variant 选择不同布局"""
        variant = decision.variant

        if variant == TemplateVariant.CONTENT_TXT_IMG:
            self._render_content_txt_img(slide, page, decision)
        elif variant == TemplateVariant.CONTENT_IMG_TXT:
            self._render_content_img_txt(slide, page, decision)
        elif variant == TemplateVariant.CONTENT_WIDE_IMG:
            self._render_content_wide_img(slide, page, decision)
        elif variant == TemplateVariant.CONTENT_TWO_COL:
            self._render_content_two_col(slide, page, decision)
        else:
            self._render_content_single_col(slide, page, decision)

    def _render_content_txt_img(self, slide, page: HTMLPageModel, decision: LayoutDecision):
        """左文右图布局"""
        text_width = Inches(5.5)
        img_left = MARGIN + text_width + Inches(0.3)
        img_width = CONTENT_WIDTH - text_width - Inches(0.3)

        # 标题
        self._add_title(slide, page.title)

        # 左侧文字
        text_lines = [tb.text for tb in page.text_blocks[:6]]
        self._add_text_body(slide, MARGIN, Inches(1.6), text_width, Inches(5.0), text_lines)

        # 右侧图片
        if page.image_blocks:
            self._add_image_from_block(
                slide, page.image_blocks[0],
                img_left, Inches(1.6), img_width, Inches(5.0),
            )

    def _render_content_img_txt(self, slide, page: HTMLPageModel, decision: LayoutDecision):
        """上图下文布局"""
        img_height = Inches(3.5)

        # 标题
        self._add_title(slide, page.title)

        # 图片（上方）
        if page.image_blocks:
            self._add_image_from_block(
                slide, page.image_blocks[0],
                MARGIN, Inches(1.4), CONTENT_WIDTH, img_height,
            )

        # 文字（下方）
        text_top = Inches(1.4) + img_height + Inches(0.3)
        text_lines = [tb.text for tb in page.text_blocks[:4]]
        self._add_text_body(
            slide, MARGIN, text_top,
            CONTENT_WIDTH, SLIDE_HEIGHT - text_top - Inches(0.6), text_lines,
        )

    def _render_content_wide_img(self, slide, page: HTMLPageModel, decision: LayoutDecision):
        """大图 + 底部文字"""
        img_height = Inches(5.5)

        # 标题
        self._add_title(slide, page.title)

        # 图片（大）
        if page.image_blocks:
            self._add_image_from_block(
                slide, page.image_blocks[0],
                MARGIN, Inches(1.2), CONTENT_WIDTH, img_height,
            )

        # 文字（底部一行）
        if page.text_blocks:
            tf = self._add_text_box(
                slide, MARGIN, Inches(1.2) + img_height + Inches(0.2),
                CONTENT_WIDTH, Inches(0.6),
            )
            p = tf.paragraphs[0]
            self._set_run(p, page.text_blocks[0].text[:300], size=14, color=COLORS["text_secondary"])

    def _render_content_two_col(self, slide, page: HTMLPageModel, decision: LayoutDecision):
        """双栏纯文本"""
        col_width = (CONTENT_WIDTH - Inches(0.3)) // 2
        items = page.text_blocks[:8]
        half = (len(items) + 1) // 2
        left_items = items[:half]
        right_items = items[half:]

        self._add_title(slide, page.title)

        # 左栏
        left_text = "\n\n".join(tb.text for tb in left_items)
        tf_l = self._add_text_box(
            slide, MARGIN, Inches(1.6), col_width, Inches(5.0),
        )
        p_l = tf_l.paragraphs[0]
        self._set_run(p_l, left_text, size=15, color=COLORS["text_primary"])

        # 右栏
        right_text = "\n\n".join(tb.text for tb in right_items)
        right_left = MARGIN + col_width + Inches(0.3)
        tf_r = self._add_text_box(
            slide, right_left, Inches(1.6), col_width, Inches(5.0),
        )
        p_r = tf_r.paragraphs[0]
        self._set_run(p_r, right_text, size=15, color=COLORS["text_primary"])

    def _render_content_single_col(self, slide, page: HTMLPageModel, decision: LayoutDecision):
        """单栏纯文本"""
        self._add_title(slide, page.title)
        text_lines = [tb.text for tb in page.text_blocks[:6]]
        self._add_text_body(slide, MARGIN, Inches(1.6), CONTENT_WIDTH, Inches(5.0), text_lines)

    # ══════════════════════════════════════════════
    # 要点页
    # ══════════════════════════════════════════════
    def _render_bullets_page(self, slide, page: HTMLPageModel, decision: LayoutDecision):
        """要点页：卡片式 / 列表式"""
        self._add_title(slide, page.title)

        items = [tb for tb in page.text_blocks if tb.is_bullet] or page.text_blocks[:6]
        item_height = Inches(0.9)
        start_top = Inches(1.8)
        gap = Inches(0.15)

        for i, tb in enumerate(items[:6]):
            top = start_top + i * (item_height + gap)
            # 卡片背景
            card = slide.shapes.add_shape(
                1, MARGIN, top, CONTENT_WIDTH, item_height,
            )
            card.fill.solid()
            card.fill.fore_color.rgb = COLORS["bg_soft"]
            card.line.color.rgb = COLORS["border"]
            card.line.width = Pt(0.5)

            # 文字
            tf = card.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            self._set_run(p, f"• {tb.text}", size=16, color=COLORS["text_primary"])

    # ══════════════════════════════════════════════
    # KPI 网格页
    # ══════════════════════════════════════════════
    def _render_kpi_grid_page(self, slide, page: HTMLPageModel, decision: LayoutDecision):
        """KPI 卡片网格"""
        self._add_title(slide, page.title)

        metrics = page.kpi_metrics[:6]
        if not metrics:
            return

        cols = min(len(metrics), 4)
        if cols == 3:
            cols = 3
        elif cols == 4:
            cols = 4
        else:
            cols = 2 if len(metrics) == 1 else min(len(metrics), 2)

        card_width = (CONTENT_WIDTH - Inches(0.3) * (cols - 1)) // cols
        card_height = Inches(3.5)
        start_top = Inches(2.0)

        for i, m in enumerate(metrics):
            col = i % cols
            row = i // cols
            left = MARGIN + col * (card_width + Inches(0.3))
            top = start_top + row * (card_height + Inches(0.2))

            # 卡片
            card = slide.shapes.add_shape(1, left, top, card_width, card_height)
            card.fill.solid()
            card.fill.fore_color.rgb = COLORS["bg_soft"]
            card.line.color.rgb = COLORS["border"]
            card.line.width = Pt(0.5)

            tf = card.text_frame
            tf.word_wrap = True

            # KPI 大数字
            p_val = tf.paragraphs[0]
            self._set_run(p_val, str(m.get("value", "")), size=36, bold=True,
                         color=COLORS["accent"], alignment=PP_ALIGN.CENTER)

            # Label
            p_label = tf.add_paragraph()
            self._set_run(p_label, m.get("label", ""), size=14,
                         color=COLORS["text_secondary"], alignment=PP_ALIGN.CENTER)

            # Change indicator
            change = m.get("change", "")
            if change:
                p_change = tf.add_paragraph()
                trend = m.get("trend", "up")
                self._set_run(p_change, change, size=16, bold=True,
                             color=COLORS["good"] if trend == "up" else COLORS["bad"],
                             alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════
    # 其他布局
    # ══════════════════════════════════════════════
    def _render_two_column_page(self, slide, page: HTMLPageModel, decision: LayoutDecision):
        self._render_content_two_col(slide, page, decision)

    def _render_three_column_page(self, slide, page: HTMLPageModel, decision: LayoutDecision):
        """三列卡片"""
        self._add_title(slide, page.title)
        items = page.text_blocks[:6]
        col_width = (CONTENT_WIDTH - Inches(0.6)) // 3
        for i, tb in enumerate(items):
            left = MARGIN + i * (col_width + Inches(0.3))
            card = slide.shapes.add_shape(1, left, Inches(1.8), col_width, Inches(4.0))
            card.fill.solid()
            card.fill.fore_color.rgb = COLORS["bg_soft"]
            card.line.color.rgb = COLORS["border"]
            card.line.width = Pt(0.5)
            tf = card.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            self._set_run(p, tb.text[:200], size=15, color=COLORS["text_primary"])

    def _render_table_page(self, slide, page: HTMLPageModel, decision: LayoutDecision):
        """表格页"""
        self._add_title(slide, page.title)

        if not page.table_data:
            return

        headers = page.table_data.get("headers", [])
        rows = page.table_data.get("rows", [])[:10]

        if not headers or not rows:
            return

        n_rows = len(rows) + 1
        n_cols = len(headers)

        table_top = Inches(2.0)
        table_height = Inches(0.4) * n_rows
        table_width = CONTENT_WIDTH

        table_shape = slide.shapes.add_table(
            n_rows, n_cols,
            MARGIN, table_top, table_width, table_height,
        )
        table = table_shape.table

        # Headers
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = h
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.bold = True
                paragraph.font.color.rgb = COLORS["white"]
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["accent"]

        # Data rows
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.text = str(val)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(12)
                    paragraph.font.color.rgb = COLORS["text_primary"]
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLORS["bg_soft"]

    def _render_image_hero_page(self, slide, page: HTMLPageModel, decision: LayoutDecision):
        """大图页"""
        self._add_title(slide, page.title)
        if page.image_blocks:
            self._add_image_from_block(
                slide, page.image_blocks[0],
                MARGIN, Inches(1.6), CONTENT_WIDTH, Inches(5.0),
            )

    def _render_image_grid_page(self, slide, page: HTMLPageModel, decision: LayoutDecision):
        """图片网格"""
        self._add_title(slide, page.title)
        images = page.image_blocks[:6]
        if not images:
            return
        cols = 3 if len(images) >= 3 else len(images)
        img_width = (CONTENT_WIDTH - Inches(0.3) * (cols - 1)) // cols
        img_height = Inches(2.5)
        for i, img in enumerate(images[:6]):
            col = i % cols
            row = i // cols
            left = MARGIN + col * (img_width + Inches(0.3))
            top = Inches(2.0) + row * (img_height + Inches(0.2))
            self._add_image_from_block(slide, img, left, top, img_width, img_height)

    def _render_stat_page(self, slide, page: HTMLPageModel, decision: LayoutDecision):
        """大数字页"""
        center_x = SLIDE_WIDTH // 2
        if page.kpi_metrics:
            m = page.kpi_metrics[0]
            tf = self._add_text_box(
                slide, center_x - Inches(4), Inches(1.8), Inches(8), Inches(2.0),
            )
            self._set_run(tf.paragraphs[0], str(m.get("value", "0")),
                         size=120, bold=True, color=COLORS["accent"], alignment=PP_ALIGN.CENTER)

            tf2 = self._add_text_box(
                slide, center_x - Inches(4), Inches(4.0), Inches(8), Inches(1.0),
            )
            self._set_run(tf2.paragraphs[0], m.get("label", page.title),
                         size=24, color=COLORS["text_primary"], alignment=PP_ALIGN.CENTER)

    def _render_thanks_page(self, slide, page: HTMLPageModel, decision: LayoutDecision):
        """致谢页"""
        center_x = SLIDE_WIDTH // 2
        center_y = SLIDE_HEIGHT // 2
        tf = self._add_text_box(
            slide, center_x - Inches(4), center_y - Inches(1.5),
            Inches(8), Inches(3.0),
        )
        self._set_run(tf.paragraphs[0], page.title or "Thank You",
                     size=72, bold=True, color=COLORS["accent"], alignment=PP_ALIGN.CENTER)

    def _render_chart_page(self, slide, page: HTMLPageModel, decision: LayoutDecision):
        """图表页（使用 python-pptx 原生 Chart）"""
        self._add_title(slide, page.title)

        if not page.chart_data:
            return

        labels = page.chart_data.get("labels", [])
        datasets = page.chart_data.get("datasets", [])

        if not labels or not datasets:
            return

        # 构建 Chart.js 风格数据
        chart_data = CategoryChartData()
        chart_data.categories = labels
        for ds in datasets:
            chart_data.add_series(ds.get("label", "Series"), ds.get("data", []))

        # 图表类型映射
        chart_type_map = {
            LayoutType.CHART_BAR: XL_CHART_TYPE.COLUMN_CLUSTERED,
            LayoutType.CHART_LINE: XL_CHART_TYPE.LINE,
            LayoutType.CHART_PIE: XL_CHART_TYPE.PIE,
        }
        chart_type = chart_type_map.get(page.layout, XL_CHART_TYPE.COLUMN_CLUSTERED)

        # 插入图表
        chart_frame = slide.shapes.add_chart(
            chart_type,
            MARGIN, Inches(1.8), CONTENT_WIDTH, Inches(5.0),
            chart_data,
        )

        chart = chart_frame.chart
        chart.has_legend = True
        chart.legend.include_in_layout = False

    def _render_fallback(self, slide, page: HTMLPageModel, idx: int, total: int):
        """兜底页：简单标题 + 文本"""
        tf = self._add_text_box(slide, MARGIN, Inches(1.0), CONTENT_WIDTH, Inches(1.0))
        self._set_run(tf.paragraphs[0], page.title or f"Slide {idx + 1}", size=24, bold=True)

        if page.text_blocks:
            lines = [tb.text[:300] for tb in page.text_blocks[:4]]
            tf2 = self._add_text_box(slide, MARGIN, Inches(2.5), CONTENT_WIDTH, Inches(4.0))
            self._set_run(tf2.paragraphs[0], "\n\n".join(lines), size=14, color=COLORS["text_primary"])

    # ══════════════════════════════════════════════
    # 工具方法
    # ══════════════════════════════════════════════

    @staticmethod
    def _strip_markdown(text: str) -> str:
        """清洗 Markdown 格式标记（**bold**, *italic*, `code`, ## heading, [link]()等）"""
        import re
        text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)    # **bold**
        text = re.sub(r'\*(.+?)\*', r'\1', text)          # *italic*
        text = re.sub(r'`(.+?)`', r'\1', text)             # `code`
        text = re.sub(r'~~(.+?)~~', r'\1', text)           # ~~strikethrough~~
        text = re.sub(r'^#{1,6}\s*', '', text, flags=re.MULTILINE)  # # headings
        text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)       # [text](url)
        text = re.sub(r'!\[([^\]]*)\]\([^\)]+\)', r'\1', text)      # ![alt](img)
        text = re.sub(r'>\s?', '', text)                 # > blockquote
        text = re.sub(r'---+\s*', '', text)              # --- separator
        text = re.sub(r'\|\s*', ' ', text)               # | table pipes → spaces
        text = re.sub(r'\[\^?\d+\]', '', text)           # [^1] footnotes
        text = re.sub(r'\s+', ' ', text).strip()         # normalize whitespace
        return text

    def _set_run(
        self, paragraph, text: str,
        size: int = 15, bold: bool = False,
        color=None, alignment=None, strip_md: bool = True,
    ):
        """安全地在段落中添加带正确字体的 Run

        用 p.add_run() 替代 p.text = ... 确保字体在 PPTX 中正确持久化。
        python-pptx 中 p.text = 创建的隐式 run 不会继承 p.font 设置。
        """
        clean = self._strip_markdown(text) if strip_md else text
        run = paragraph.add_run()
        run.text = clean
        run.font.name = FONT_FAMILY
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        if alignment:
            paragraph.alignment = alignment
        return run

    def _add_title(self, slide, title: str):
        """添加页面标题"""
        tf = self._add_text_box(slide, MARGIN, Inches(0.4), CONTENT_WIDTH, Inches(0.8))
        p = tf.paragraphs[0]
        self._set_run(p, title[:100], size=28, bold=True, color=COLORS["text_primary"])

    def _add_text_body(self, slide, left, top, width, height, lines: List[str]):
        """添加正文文本框"""
        tf = self._add_text_box(slide, left, top, width, height)
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
                p.space_before = Pt(8)
            self._set_run(p, line[:500], size=15, color=COLORS["text_primary"])

    def _add_image_from_block(
        self, slide, img_block: HTMLImageBlock,
        left, top, width, height,
    ):
        """从 HTMLImageBlock 添加图片"""
        url = img_block.url or ""
        if not url:
            return

        try:
            if url.startswith("data:image"):
                # Base64 图片 — 修复 padding 问题
                header, b64_data = url.split(",", 1)
                # 确保 padding 正确（base64 长度必须是 4 的倍数）
                b64_clean = b64_data.strip().replace('\n', '').replace('\r', '').replace(' ', '')
                missing_padding = len(b64_clean) % 4
                if missing_padding:
                    b64_clean += '=' * (4 - missing_padding)
                img_bytes = base64.b64decode(b64_clean, validate=True)
                img_stream = io.BytesIO(img_bytes)
                slide.shapes.add_picture(img_stream, left, top, width, height)
            elif url.startswith("file:///"):
                # 本地文件
                file_path = url.replace("file:///", "")
                slide.shapes.add_picture(file_path, left, top, width, height)
            else:
                logger.debug(f"Skipping non-base64 image: {url[:50]}...")
        except Exception as e:
            logger.warning(f"Failed to add image: {e}")

    def _add_text_box(self, slide, left, top, width, height):
        """添加文本框"""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        txBox.text_frame.word_wrap = True
        return txBox.text_frame

    def _add_rounded_rect(self, slide, left, top, width, height, fill_color):
        """添加圆角矩形"""
        shape = slide.shapes.add_shape(
            5,  # MSO_SHAPE.ROUNDED_RECTANGLE
            left, top, width, height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        # 设置文字样式
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS["white"]
        p.font.name = FONT_FAMILY
        p.alignment = PP_ALIGN.CENTER
        return tf

    def _add_footer(self, slide, page: HTMLPageModel, idx: int, total: int):
        """添加页脚"""
        tf = self._add_text_box(
            slide, MARGIN, SLIDE_HEIGHT - Inches(0.5),
            CONTENT_WIDTH, Inches(0.4),
        )
        self._set_run(tf.paragraphs[0], f"{page.kicker or ''}  ·  {idx + 1}/{total}",
                     size=10, color=COLORS["text_tertiary"], alignment=PP_ALIGN.RIGHT)


# Module-level convenience
def render_pptx_native(
    pages: List[HTMLPageModel],
    output_path: str,
    title: str = "Market Insight Report",
    template_path: Optional[str] = None,
) -> str:
    """便捷函数：原生 PPTX 渲染"""
    renderer = PPTXNativeRenderer(template_path=template_path)
    return renderer.render(pages, output_path, title)

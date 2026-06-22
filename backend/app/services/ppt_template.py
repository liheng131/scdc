"""
PPT 母版模板加载与填充服务

核心职责：
1. 扫描 `backend/app/templates/ppt/` 目录下的 .pptx 母版
2. 解析每个母版的 slide_master 与 slide_layouts
3. 提供 list_templates() / get_template() / fill_template() 三个对外接口
4. 内部缓存已加载的模板对象，避免每次导出重新解析

设计原则：
- 模板 ID 来自 MANIFEST.json，未配置时回退到文件名（template{N}）
- layout 选取策略：先按 layouts_hint 名称匹配，匹配不到时按布局功能（cover/section/content/picture/summary）选最接近的
- 所有数据填充通过 add_slide(layout) + placeholder.text_frame 操作，不破坏母版的背景/字体/配色
- 长内容自动分页；配图按比例缩放
"""

import io
import base64
import json
import logging
import os
import re
from dataclasses import dataclass, field
from typing import Any, Dict, List, Optional, Tuple
from datetime import datetime

from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

from app.models.report import Report
from app.services.report_page_model import (
    ReportPageModel, PageModel, TextBlock, ImageBlock, TableBlock,
    USABLE_WIDTH, USABLE_HEIGHT,
    TEXT_AREA_MAX_HEIGHT, IMAGE_AREA_MAX_HEIGHT,
)

# ── 图片最小尺寸常量 ──
MIN_IMG_WIDTH = Inches(1.5)   # 约 144px @ 96dpi, 保证可读
MIN_IMG_HEIGHT = Inches(1.0)  # 约 96px @ 96dpi, 保证可读

logger = logging.getLogger(__name__)

# 模板根目录（可被环境变量覆盖，便于测试）
TEMPLATES_DIR = os.environ.get(
    "PPT_TEMPLATES_DIR",
    os.path.join(os.path.dirname(os.path.dirname(__file__)), "templates", "ppt"),
)

# 单幻灯片正文承载上限（字符），超过则自动分页
MAX_CHARS_PER_CONTENT_SLIDE = 500
MAX_CHARS_PER_PARAGRAPH = 800

# 统一字体（中英文）
CHINESE_FONT = "Microsoft YaHei"
LATIN_FONT = "Calibri"

# 默认文字色（深灰，浅色背景时使用）
DEFAULT_TEXT_COLOR = RGBColor(0x33, 0x33, 0x33)
DEFAULT_TITLE_COLOR = RGBColor(0x1F, 0x3A, 0x5F)  # 深蓝标题色
WHITE_COLOR = RGBColor(0xFF, 0xFF, 0xFF)


def _set_run_font(run, name=CHINESE_FONT, size=None, bold=None, color=None, italic=None):
    """统一设置 run 字体（同时设置西文 + 东亚字体）。
    - name: 中文字体名（同时作为东亚字体）
    - size: Pt 字号
    - bold/italic: 布尔
    - color: RGBColor
    """
    # 设置西文字体
    run.font.name = LATIN_FONT
    # 通过 XML 设置东亚字体（这是关键：让中文也用指定字体）
    rPr = run._r.get_or_add_rPr()
    # 移除现有的 latin/ea/cs 字体设置
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        for el in rPr.findall(qn(tag)):
            rPr.remove(el)
    # 添加新的 latin（西文）
    latin = etree.SubElement(rPr, qn('a:latin'))
    latin.set('typeface', LATIN_FONT)
    # 添加东亚字体（中文）
    ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', name)
    # 添加复杂脚本字体
    cs = etree.SubElement(rPr, qn('a:cs'))
    cs.set('typeface', LATIN_FONT)
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color


def _calc_adaptive_font_size(text_length: int) -> int:
    """根据内容长度选择字号,确保不溢出可视区域"""
    if text_length < 200:
        return 18  # 短内容用大字号
    elif text_length < 400:
        return 16
    elif text_length < 600:
        return 14
    else:
        return 13  # 长内容用小字号


def _calculate_luminance(rgb_hex: str) -> float:
    """计算背景色亮度（0-255）"""
    if not rgb_hex:
        return 255.0
    hex_str = rgb_hex.lstrip('#')
    if len(hex_str) < 6:
        return 255.0
    try:
        r = int(hex_str[0:2], 16)
        g = int(hex_str[2:4], 16)
        b = int(hex_str[4:6], 16)
        # 标准亮度公式
        return (r * 299 + g * 587 + b * 114) / 1000
    except Exception:
        return 255.0


def _pick_text_color(bg_luminance: float) -> RGBColor:
    """根据背景亮度选文字色（深色背景用白色，浅色背景用深色）"""
    if bg_luminance < 128:
        return WHITE_COLOR
    return DEFAULT_TEXT_COLOR


def _make_unique(name: str, suffix: int) -> str:
    """生成唯一的 zip entry 名称：name -> name_dup{suffix}"""
    if "." in name:
        base, ext = name.rsplit(".", 1)
        return f"{base}_dup{suffix}.{ext}"
    return f"{name}_dup{suffix}"


# ─────────────────────────── 数据结构 ───────────────────────────


@dataclass
class TemplateInfo:
    """模板元数据（API 返回）"""
    id: str
    name: str
    description: str
    file: str
    layouts_count: int = 0


@dataclass
class TemplateContext:
    """已加载的模板对象（包含 Presentation 实例和解析后的 layout 映射）"""
    info: TemplateInfo
    file_path: str
    presentation: Presentation
    cover_layouts: List[Any] = field(default_factory=list)      # 封面
    section_layouts: List[Any] = field(default_factory=list)    # 章节封面
    content_layouts: List[Any] = field(default_factory=list)    # 正文
    picture_layouts: List[Any] = field(default_factory=list)    # 配图
    summary_layouts: List[Any] = field(default_factory=list)    # 总结

    def pick_cover(self, idx: int = 0):
        return self.cover_layouts[idx] if self.cover_layouts else (
            self.section_layouts[0] if self.section_layouts else
            (self.content_layouts[0] if self.content_layouts else self.presentation.slide_layouts[0])
        )

    def pick_section(self, idx: int = 0):
        return self.section_layouts[idx % len(self.section_layouts)] if self.section_layouts else self.pick_cover()

    def pick_content(self, idx: int = 0):
        return self.content_layouts[idx % len(self.content_layouts)] if self.content_layouts else self.pick_cover()

    def pick_picture(self, idx: int = 0):
        return self.picture_layouts[idx % len(self.picture_layouts)] if self.picture_layouts else self.pick_content()

    def pick_summary(self, idx: int = 0):
        return self.summary_layouts[idx % len(self.summary_layouts)] if self.summary_layouts else self.pick_section()


# ─────────────────────────── 模板服务 ───────────────────────────


class PPTTemplateService:
    """PPT 母版模板加载与填充服务（单例）"""

    _instance: Optional["PPTTemplateService"] = None
    _templates: Dict[str, TemplateContext] = {}
    _manifest: Dict[str, Any] = {}

    def __new__(cls):
        if cls._instance is None:
            cls._instance = super().__new__(cls)
        return cls._instance

    # ── 公开 API ─────────────────────────────────────────────

    def list_templates(self) -> List[TemplateInfo]:
        """返回所有已加载模板的元数据列表（供前端拉取）"""
        self._ensure_loaded()
        return [ctx.info for ctx in self._templates.values()]

    def get_template(self, template_id: str) -> Optional[TemplateContext]:
        """按 ID 返回缓存的模板上下文，找不到时返回 None"""
        self._ensure_loaded()
        return self._templates.get(template_id)

    def fill_template(
        self,
        template_id: str,
        report: Report,
        chart_images: Optional[List[Dict[str, str]]] = None,
    ) -> bytes:
        """核心入口：用指定模板生成填充后的 PPTX 字节流"""
        self._ensure_loaded()
        ctx = self._templates.get(template_id)
        if not ctx:
            logger.warning("Template '%s' not found, falling back to first available", template_id)
            ctx = next(iter(self._templates.values()), None)
        if not ctx:
            # 无任何模板可用，回退到内置默认（与旧 generate_pptx 行为一致）
            return self._fallback_pptx(report, chart_images)

        # 每次生成都从模板文件重新加载 Presentation,确保干净状态
        prs = Presentation(ctx.file_path)

        # ── 关键步骤: 清空模板原 slide,只保留 slide_masters / slide_layouts ──
        # 原因:模板是用户手工制作的"设计母版",里面通常有 5 张示例 slide。
        # 如果不清空,新追加的报告 slide 会接在示例 slide 后面,PowerPoint
        # 默认从第 1 张展示,导致用户看到模板原内容("月度市场动态(国内)")
        # 而不是报告内容。
        try:
            original_slide_count = len(prs.slides)
            xml_slides = prs.slides._sldIdLst
            # 必须先复制成 list 再 remove,否则会边遍历边修改
            slides_to_remove = list(xml_slides)
            for slide_id in slides_to_remove:
                xml_slides.remove(slide_id)
            logger.info(
                "[PPT-TEMPLATE] Cleared %d original slides from template '%s', "
                "keeping only slide_masters/slide_layouts for design consistency",
                original_slide_count, template_id,
            )
        except Exception as clear_e:
            logger.warning(
                "[PPT-TEMPLATE] Failed to clear original slides from template '%s': "
                "%s. Generated PPT may contain template example slides.",
                template_id, clear_e, exc_info=True,
            )

        # 1. 解析报告 Markdown 为结构化大纲
        outline = self._parse_outline(report)

        # 2. 封面 slide
        self._add_cover_slide(prs, ctx, report)

        # 3. 目录 slide（自动列出所有 ## 章节）
        section_titles = [s["title"] for s in outline["sections"]]
        if section_titles:
            self._add_toc_slide(prs, ctx, report.title, section_titles)

        # 4. 每个章节：章节封面 + 正文（自动分页）
        for sec_idx, section in enumerate(outline["sections"], start=1):
            # 4.1 章节封面
            self._add_section_cover_slide(prs, ctx, sec_idx, section)

            # 4.2 正文（按 MAX_CHARS_PER_CONTENT_SLIDE 自动分页）
            for page_idx, chunk in enumerate(self._split_section_content(section["content"])):
                self._add_content_slide(prs, ctx, section["title"], chunk, page_idx=page_idx)

            # 4.3 该章节末尾插入配图（按 position 匹配）
            sec_charts = [c for c in (chart_images or []) if c.get("position") == sec_idx]
            for ci in sec_charts:
                self._add_picture_slide(prs, ctx, ci)

        # 5. 总结 slide
        self._add_summary_slide(prs, ctx, report, outline["sections"])

        # 6. 序列化
        buf = io.BytesIO()
        prs.save(buf)
        raw = buf.getvalue()
        # 7. 后处理：去重 zip 内部重名 part（用户原始模板的多 master 结构会触发此问题）
        return self._dedupe_zip_parts(raw)

    # ── 新接口：基于 ReportPageModel 的统一渲染 ──────────────

    def fill_template_from_model(
        self,
        template_id: str,
        model: ReportPageModel,
    ) -> bytes:
        """基于统一的 ReportPageModel 渲染 PPTX 文件。

        与 fill_template() 的区别：
        - 输入是已解析好的 ReportPageModel（而非 Report 原始对象）
        - 不再内部解析 Markdown → 分页，而是直接遍历 model.pages
        - 保证 PPT 第 N 页 = PDF 第 N 页 = Word 第 N 页

        Args:
            template_id: 模板 ID
            model: ReportPageModel 实例

        Returns:
            PPTX 字节流
        """
        self._ensure_loaded()
        ctx = self._templates.get(template_id)
        if not ctx:
            logger.warning("Template '%s' not found, falling back to first available", template_id)
            ctx = next(iter(self._templates.values()), None)
        if not ctx:
            return self._fallback_pptx_from_model(model)

        prs = Presentation(ctx.file_path)

        # 清空模板原 slide
        try:
            xml_slides = prs.slides._sldIdLst
            for slide_id in list(xml_slides):
                xml_slides.remove(slide_id)
        except Exception as clear_e:
            logger.warning("[PPT-TEMPLATE] Failed to clear original slides: %s", clear_e, exc_info=True)

        # 遍历 model.pages 逐页渲染
        for page in model.pages:
            self._render_page(prs, ctx, page)

        buf = io.BytesIO()
        prs.save(buf)
        raw = buf.getvalue()
        return self._dedupe_zip_parts(raw)

    def _render_page(self, prs: Presentation, ctx, page: PageModel):
        """根据 page_type 分发到对应的渲染方法"""
        page_type = page.page_type
        if page_type == "cover":
            self._render_cover_from_model(prs, ctx, page)
        elif page_type == "toc":
            self._render_toc_from_model(prs, ctx, page)
        elif page_type == "section":
            self._render_section_from_model(prs, ctx, page)
        elif page_type == "content":
            self._render_content_from_model(prs, ctx, page)
        elif page_type == "picture":
            self._render_picture_from_model(prs, ctx, page)
        elif page_type == "summary":
            self._render_summary_from_model(prs, ctx, page)
        else:
            logger.warning("Unknown page_type '%s', treating as content", page_type)
            self._render_content_from_model(prs, ctx, page)

    @staticmethod
    def _clear_empty_placeholders(slide):
        """清除所有含"单击此处添加文本"/"Click to add text"的占位符"""
        placeholder_texts = {
            "单击此处添加文本",
            "Click to add text",
            "单击此处添加标题",
            "Click to add title",
            "单击此处添加副标题",
            "Click to add subtitle",
        }
        for ph in list(slide.placeholders):
            try:
                tf = ph.text_frame
                if tf.text and tf.text.strip() in placeholder_texts:
                    tf.clear()
                    if tf.paragraphs:
                        tf.paragraphs[0].text = ""
            except Exception:
                pass

    def _render_cover_from_model(self, prs: Presentation, ctx, page: PageModel):
        layout = ctx.pick_cover(0)
        slide = prs.slides.add_slide(layout)

        # 标题
        title_block = next((tb for tb in page.text_blocks if tb.style == "title"), None)
        if title_block:
            self._set_placeholder_text(
                slide,
                lambda i, t, n: t in (3, 1) or "标题" in n,
                title_block.text,
                font_size=title_block.font_size,
                bold=title_block.bold,
                color=self._hex_to_rgb(title_block.color),
                alignment=PP_ALIGN.CENTER,
            )

        # 副标题
        sub_block = next((tb for tb in page.text_blocks if tb.style == "subtitle"), None)
        if sub_block:
            self._set_placeholder_text(
                slide,
                lambda i, t, n: t == 4 or "副标题" in n,
                sub_block.text,
                font_size=sub_block.font_size,
                color=self._hex_to_rgb(sub_block.color),
                alignment=PP_ALIGN.CENTER,
            )

        # 清空未填充的占位符
        self._clear_empty_placeholders(slide)

    def _render_toc_from_model(self, prs: Presentation, ctx, page: PageModel):
        layout = ctx.pick_content(0) if ctx.content_layouts else ctx.pick_cover(0)
        slide = prs.slides.add_slide(layout)

        # 标题 "目录"
        title_block = next((tb for tb in page.text_blocks if tb.style == "title"), None)
        if title_block:
            self._set_placeholder_text(
                slide,
                lambda i, t, n: t in (1, 3) or "标题" in n,
                title_block.text,
                font_size=title_block.font_size,
                bold=title_block.bold,
                color=self._hex_to_rgb(title_block.color),
            )

        # 目录项
        body_ph = self._find_placeholder(slide, type=7) or self._find_placeholder(slide, type=2)
        if body_ph is None:
            body_ph = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.0), Inches(5.5))
        try:
            body_ph.left = Inches(0.8)
            body_ph.top = Inches(1.6)
            body_ph.width = Inches(11.7)
            body_ph.height = Inches(5.5)
        except Exception:
            pass
        tf = body_ph.text_frame
        tf.word_wrap = True
        tf.clear()

        body_blocks = [tb for tb in page.text_blocks if tb.style != "title"]
        for i, tb in enumerate(body_blocks):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = tb.text
            p.level = 0
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(8)
            p.space_after = Pt(8)
            run = p.runs[0] if p.runs else p.add_run()
            run.text = p.text
            _set_run_font(run, size=tb.font_size, color=self._hex_to_rgb(tb.color))

        # 清空未填充的占位符
        self._clear_empty_placeholders(slide)

    def _render_section_from_model(self, prs: Presentation, ctx, page: PageModel):
        layout = ctx.pick_section(page.section_number - 1)
        slide = prs.slides.add_slide(layout)

        title_block = next((tb for tb in page.text_blocks if tb.style == "title"), None)
        if title_block:
            title_set = self._set_placeholder_text(
                slide,
                lambda i, t, n: t in (1, 3) or "标题" in n,
                title_block.text,
                font_size=title_block.font_size,
                bold=title_block.bold,
                color=self._hex_to_rgb(title_block.color),
            )
            if not title_set:
                tx = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(2.0))
                tf = tx.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = title_block.text
                p.alignment = PP_ALIGN.CENTER
                run = p.runs[0] if p.runs else p.add_run()
                run.text = title_block.text
                _set_run_font(run, size=title_block.font_size, bold=title_block.bold,
                              color=self._hex_to_rgb(title_block.color))

        # 清空未填充的占位符
        self._clear_empty_placeholders(slide)

    def _render_content_from_model(self, prs: Presentation, ctx, page: PageModel):
        layout = ctx.pick_content(0)
        slide = prs.slides.add_slide(layout)

        # 标题
        title_block = next((tb for tb in page.text_blocks if tb.style == "title"), None)
        if title_block:
            title_set = self._set_placeholder_text(
                slide,
                lambda i, t, n: t in (1, 3) or "标题" in n,
                title_block.text,
                font_size=title_block.font_size,
                bold=title_block.bold,
                color=self._hex_to_rgb(title_block.color),
            )
            if not title_set:
                tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0))
                tf = tx.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = title_block.text
                p.alignment = PP_ALIGN.LEFT
                run = p.runs[0] if p.runs else p.add_run()
                run.text = title_block.text
                _set_run_font(run, size=title_block.font_size, bold=title_block.bold,
                              color=self._hex_to_rgb(title_block.color))

        # 正文块（排除标题）
        body_blocks = [tb for tb in page.text_blocks if tb.style != "title"]

        if page.layout_hint == "text_top" and page.images:
            # 上文字下图片
            self._render_text_top_layout(slide, body_blocks, page.images)
        else:
            # 纯文字
            self._render_text_only_layout(slide, body_blocks)

        # 渲染表格（在文字/图片区域下方）
        if page.tables:
            # 动态计算表格起始位置
            if page.images:
                # 有图片：放在图片区下方
                total_chars = sum(len(tb.text) for tb in body_blocks)
                chars_per_line = max(1, int(USABLE_WIDTH / 0.167))
                line_count = max(1, (total_chars + chars_per_line - 1) // chars_per_line)
                text_bottom = Inches(1.5) + Inches(min(float(TEXT_AREA_MAX_HEIGHT), 2.8 + line_count * 0.25))
                table_top = text_bottom + Inches(IMAGE_AREA_MAX_HEIGHT) + Inches(0.3)
            else:
                # 纯文字：放在文字区下方
                table_top = Inches(4.5)
            for tbl in page.tables:
                self._render_table(slide, tbl, table_top)
                table_top += Inches(1.5)  # 每个表格占 1.5 inch

    def _render_table(self, slide, table_block: "TableBlock", top: int):
        """在幻灯片中渲染一个表格"""
        from pptx.util import Pt as PptxPt
        rows = len(table_block.rows) + 1  # +1 for header
        cols = len(table_block.headers) if table_block.headers else (len(table_block.rows[0]) if table_block.rows else 2)
        if cols == 0 or rows == 1:
            return

        table_width = Inches(USABLE_WIDTH)
        table_height = Inches(min(rows * 0.35, 2.0))
        left = Inches(0.5)

        tbl_shape = slide.shapes.add_table(rows, cols, left, top, table_width, table_height)
        table = tbl_shape.table

        # 设置列宽均匀分布
        col_w = Emu(int(table_width / cols))
        for ci in range(cols):
            table.columns[ci].width = col_w

        # 填充表头
        for ci, header in enumerate(table_block.headers):
            cell = table.cell(0, ci)
            cell.text = header
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    _set_run_font(run, size=10, bold=True, color=self._hex_to_rgb("#FFFFFF"))
            # 表头背景色
            tcPr = cell._tc.get_or_add_tcPr()
            solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
            srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
            srgbClr.set('val', '1F3A5F')

        # 填充数据行
        for ri, row_data in enumerate(table_block.rows):
            for ci, val in enumerate(row_data):
                if ci < cols:
                    cell = table.cell(ri + 1, ci)
                    cell.text = str(val)
                    for p in cell.text_frame.paragraphs:
                        p.alignment = PP_ALIGN.CENTER
                        for run in p.runs:
                            _set_run_font(run, size=9, bold=False, color=DEFAULT_TEXT_COLOR)
                    # 交替行背景色
                    if ri % 2 == 1:
                        tcPr = cell._tc.get_or_add_tcPr()
                        solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
                        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
                        srgbClr.set('val', 'E8EDF2')

    def _render_text_top_layout(self, slide, text_blocks: List[TextBlock], images: List[ImageBlock]):
        """上文字下图片布局 - 图片位置根据文字高度动态计算，防止遮挡"""
        # 文字区
        text_top = Inches(1.5)
        text_height = Inches(min(TEXT_AREA_MAX_HEIGHT, 2.8))
        body_ph = self._find_placeholder(slide, type=7) or self._find_placeholder(slide, type=2)
        if body_ph is not None:
            try:
                body_ph.left = Inches(0.5)
                body_ph.top = text_top
                body_ph.width = Inches(USABLE_WIDTH)
                body_ph.height = text_height
            except Exception:
                pass
            tf = body_ph.text_frame
            tf.word_wrap = True
            tf.clear()
        else:
            tx = slide.shapes.add_textbox(
                Inches(0.5), text_top,
                Inches(USABLE_WIDTH), text_height,
            )
            tf = tx.text_frame
            tf.word_wrap = True
            tf.clear()

        if text_blocks:
            self._fill_text_blocks(tf, text_blocks)

        # 图片区 —— 动态计算起始位置，防止遮挡文字
        if images:
            # 估算文字实际占用的高度（每行约 18pt * 行数，使用 18pt 行高近似）
            total_chars = sum(len(tb.text) for tb in text_blocks)
            chars_per_line = max(1, int(USABLE_WIDTH / 0.167))  # 12pt 中文 ≈ 12/72 = 0.167in
            line_count = max(1, (total_chars + chars_per_line - 1) // chars_per_line)
            estimated_text_height = Inches(min(
                float(TEXT_AREA_MAX_HEIGHT),
                text_height / Inches(1) + line_count * 0.25,
            ))
            # estimated_text_height 已经是 Emu 值，直接相加即可
            img_top = text_top + estimated_text_height + Inches(0.2)

            # 确保图片区不超出页面底部
            slide_bottom = Inches(7.0)  # 标准 16:9 可用底部
            max_img_area = slide_bottom - img_top
            img_area_height = min(Inches(IMAGE_AREA_MAX_HEIGHT - 0.5), max_img_area)
            if img_area_height < Inches(0.5):
                img_area_height = Inches(0.5)

            img_area_left = Inches(0.5)
            img_area_width = Inches(USABLE_WIDTH)

            n = len(images)
            gap = Inches(0.2)
            each_width = Emu(int((img_area_width - gap * max(0, n - 1)) / n))

            for i, img in enumerate(images):
                try:
                    img_blob = io.BytesIO(base64.b64decode(img.base64.strip()))
                    from PIL import Image as PILImage
                    # 验证图片有效性，避免产生 0x0 空形状
                    try:
                        with PILImage.open(img_blob) as pim:
                            w, h = pim.size
                            if w < 10 or h < 10:
                                raise ValueError(f"image too small: {w}x{h}")
                            aspect = h / w if w else 0.75
                    except Exception as _pe:
                        logger.warning("PIL cannot open image '%s': %s", img.alt, _pe)
                        continue
                    img_blob.seek(0)
                    pic_w = each_width
                    pic_h = Emu(int(pic_w * aspect))

                    # 限制最大高度
                    max_h = min(Inches(2.5), img_area_height)
                    if pic_h > max_h:
                        pic_h = max_h
                        pic_w = Emu(int(pic_h / aspect))

                    # 保证最小尺寸，防止图片太小看不清
                    if pic_w < MIN_IMG_WIDTH:
                        pic_w = MIN_IMG_WIDTH
                        pic_h = Emu(int(pic_w * aspect))
                    if pic_h < MIN_IMG_HEIGHT:
                        pic_h = MIN_IMG_HEIGHT
                        pic_w = Emu(int(pic_h / aspect))

                    # 确保不超出图片区域
                    if pic_w > img_area_width:
                        pic_w = img_area_width
                        pic_h = Emu(int(pic_w * aspect))
                    if pic_h > img_area_height:
                        pic_h = img_area_height
                        pic_w = Emu(int(pic_h / aspect))

                    v_offset = Emu(int((img_area_height - pic_h) / 2))
                    left = img_area_left + Emu(int((each_width + gap) * i))
                    slide.shapes.add_picture(img_blob, left, img_top + v_offset,
                                             width=pic_w, height=pic_h)
                except Exception as e:
                    logger.warning("Failed to embed image '%s': %s", img.alt, e)

    def _render_text_only_layout(self, slide, text_blocks: List[TextBlock]):
        """纯文字布局"""
        if not text_blocks:
            return
        body_ph = self._find_placeholder(slide, type=7) or self._find_placeholder(slide, type=2)
        if body_ph is None:
            body_ph = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5), Inches(USABLE_WIDTH), Inches(5.5),
            )
        try:
            body_ph.left = Inches(0.5)
            body_ph.top = Inches(1.5)
            body_ph.width = Inches(USABLE_WIDTH)
            body_ph.height = Inches(5.5)
        except Exception:
            pass
        tf = body_ph.text_frame
        tf.word_wrap = True
        tf.clear()
        self._fill_text_blocks(tf, text_blocks)

        # 清空未填充的占位符
        self._clear_empty_placeholders(slide)

    def _render_picture_from_model(self, prs: Presentation, ctx, page: PageModel):
        layout = ctx.pick_picture(0)
        slide = prs.slides.add_slide(layout)

        # 标题
        title_block = next((tb for tb in page.text_blocks if tb.style == "title"), None)
        if title_block:
            self._set_placeholder_text(
                slide,
                lambda i, t, n: t in (1, 3) or "标题" in n,
                title_block.text,
                font_size=title_block.font_size,
                bold=title_block.bold,
            )

        # 图片
        if page.images:
            img = page.images[0]
            pic_added = False
            for ph in slide.placeholders:
                if ph.placeholder_format.type == 18:
                    try:
                        img_blob = io.BytesIO(base64.b64decode(img.base64.strip()))
                        ph.insert_picture(img_blob)
                        pic_added = True
                    except Exception as e:
                        logger.warning("insert_picture into placeholder failed: %s", e)
                    break

            if not pic_added:
                try:
                    img_blob = io.BytesIO(base64.b64decode(img.base64.strip()))
                    slide_w = prs.slide_width
                    slide_h = prs.slide_height
                    max_w = int(slide_w * img.width_ratio)
                    max_h = int(slide_h * 0.65)
                    from PIL import Image as PILImage
                    try:
                        pil = PILImage.open(img_blob)
                        iw, ih = pil.size
                        scale = min(max_w / iw, max_h / ih, 1.0)
                        slide.shapes.add_picture(
                            img_blob,
                            int((slide_w - iw * scale) / 2),
                            Inches(1.5),
                            width=int(iw * scale),
                            height=int(ih * scale),
                        )
                    except Exception:
                        slide.shapes.add_picture(img_blob, Inches(1), Inches(1.5), width=max_w)
                except Exception as e:
                    logger.warning("Failed to add picture to slide: %s", e)

        # 多图（从第 2 张开始，第 1 张已在上面单独处理）
        if len(page.images) > 1:
            slide_w = prs.slide_width
            gap = Inches(0.2)
            n = len(page.images) - 1  # 剩余图片数
            each_w = int((slide_w - gap * (n + 1)) / n)
            img_top = Inches(5.0)
            for i, img in enumerate(page.images[1:], start=1):
                try:
                    img_blob = io.BytesIO(base64.b64decode(img.base64.strip()))
                    left = gap + (each_w + gap) * i
                    # 保证最小尺寸
                    final_w = max(each_w, MIN_IMG_WIDTH)
                    final_w = min(final_w, slide_w - left - gap)
                    slide.shapes.add_picture(img_blob, int(left), img_top, width=final_w)
                except Exception as e:
                    logger.warning("Failed to add extra picture: %s", e)

        # 清空未填充的占位符
        self._clear_empty_placeholders(slide)

    def _render_summary_from_model(self, prs: Presentation, ctx, page: PageModel):
        layout = ctx.pick_summary(0)
        slide = prs.slides.add_slide(layout)

        title_block = next((tb for tb in page.text_blocks if tb.style == "title"), None)
        if title_block:
            self._set_placeholder_text(
                slide,
                lambda i, t, n: t in (1, 3) or "标题" in n,
                title_block.text,
                font_size=title_block.font_size,
                bold=title_block.bold,
            )

        # 移除模板 BODY/OBJECT 占位符
        placeholders_to_remove = []
        for ph in list(slide.placeholders):
            try:
                if ph.placeholder_format.type in (2, 7):
                    placeholders_to_remove.append(ph)
            except Exception:
                pass
        for ph in placeholders_to_remove:
            try:
                sp = ph._element
                sp.getparent().remove(sp)
            except Exception:
                pass

        body_blocks = [tb for tb in page.text_blocks if tb.style != "title"]
        if body_blocks:
            tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9.0), Inches(5.5))
            tf = tx.text_frame
            tf.word_wrap = True
            tf.clear()
            self._fill_text_blocks(tf, body_blocks)

        # 清空未填充的占位符
        self._clear_empty_placeholders(slide)

    def _fill_text_blocks(self, tf, text_blocks: List[TextBlock]):
        """将 TextBlock 列表填充到 text_frame"""
        for i, tb in enumerate(text_blocks):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = tb.text
            p.space_after = Pt(4)
            if tb.alignment == "center":
                p.alignment = PP_ALIGN.CENTER
            elif tb.alignment == "right":
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT
            run = p.runs[0] if p.runs else p.add_run()
            run.text = tb.text
            _set_run_font(run, size=tb.font_size, bold=tb.bold,
                          color=self._hex_to_rgb(tb.color))

    @staticmethod
    def _hex_to_rgb(hex_color: str) -> RGBColor:
        """#RRGGBB → RGBColor"""
        hex_str = hex_color.lstrip('#')
        if len(hex_str) >= 6:
            return RGBColor(
                int(hex_str[0:2], 16),
                int(hex_str[2:4], 16),
                int(hex_str[4:6], 16),
            )
        return DEFAULT_TEXT_COLOR

    def _fallback_pptx_from_model(self, model: ReportPageModel) -> bytes:
        """无模板时的兜底渲染"""
        from pptx import Presentation as DefaultPresentation
        prs = DefaultPresentation()
        blank_layout = prs.slide_layouts[6]  # blank

        for page in model.pages:
            slide = prs.slides.add_slide(blank_layout)
            # 标题
            if page.title:
                tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.8))
                tf = tx.text_frame
                p = tf.paragraphs[0]
                p.text = page.title
                p.font.size = Pt(24)
                p.font.bold = True

            # 正文
            body_blocks = [tb for tb in page.text_blocks if tb.style != "title"]
            if body_blocks:
                tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(5.0))
                tf = tx.text_frame
                tf.word_wrap = True
                for i, tb in enumerate(body_blocks):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = tb.text
                    p.font.size = Pt(tb.font_size)

            # 图片
            if page.images:
                try:
                    img = page.images[0]
                    img_blob = io.BytesIO(base64.b64decode(img.base64.strip()))
                    slide.shapes.add_picture(img_blob, Inches(1), Inches(4.5), width=Inches(8))
                except Exception:
                    pass

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    # ── 内部：加载与解析 ─────────────────────────────────────

    def _ensure_loaded(self):
        """确保模板已扫描并缓存（首次访问时执行）"""
        if self._templates:
            return
        self._load_manifest()
        self._scan_templates()

    def _load_manifest(self):
        manifest_path = os.path.join(TEMPLATES_DIR, "MANIFEST.json")
        if os.path.exists(manifest_path):
            try:
                with open(manifest_path, "r", encoding="utf-8") as f:
                    self._manifest = json.load(f)
            except Exception as e:
                logger.warning("Failed to load MANIFEST.json: %s, falling back to file-based discovery", e)
                self._manifest = {}
        else:
            self._manifest = {}

    def _scan_templates(self):
        """扫描目录，加载所有 .pptx 母版并按名称归类 layout"""
        if not os.path.isdir(TEMPLATES_DIR):
            logger.warning("PPT templates dir does not exist: %s", TEMPLATES_DIR)
            return

        manifest_templates = {t["id"]: t for t in self._manifest.get("templates", [])}
        # 用绝对路径作为去重 key（同一个文件不应被加载两次）
        registered_files: set = set()

        # 先按 MANIFEST 顺序加载
        for tid, meta in manifest_templates.items():
            file_path = os.path.abspath(os.path.join(TEMPLATES_DIR, meta["file"]))
            if os.path.exists(file_path):
                self._register_template(tid, meta, file_path)
                registered_files.add(file_path)

        # 再扫描目录中未在 MANIFEST 里的 .pptx
        for fname in sorted(os.listdir(TEMPLATES_DIR)):
            if not fname.lower().endswith(".pptx"):
                continue
            file_path = os.path.abspath(os.path.join(TEMPLATES_DIR, fname))
            if file_path in registered_files:
                continue
            stem = os.path.splitext(fname)[0]
            tid = stem  # template1/template2/template3
            self._register_template(
                tid,
                {
                    "id": tid,
                    "file": fname,
                    "name": stem,
                    "description": "未在 MANIFEST 中注册的模板",
                    "layouts_hint": {},
                },
                file_path,
            )
            registered_files.add(file_path)

        logger.info("PPTTemplateService: loaded %d template(s)", len(self._templates))

    def _register_template(self, tid: str, meta: Dict[str, Any], file_path: str):
        # 优先加载去重后的副本（源文件若有重复 slide layout 名，python-pptx 保存时 zipfile 会冲突）
        usable_path = self._sanitize_template(file_path)
        try:
            prs = Presentation(usable_path)
        except Exception as e:
            logger.error("Failed to open template %s: %s", usable_path, e)
            return

        ctx = TemplateContext(
            info=TemplateInfo(
                id=meta["id"],
                name=meta.get("name", tid),
                description=meta.get("description", ""),
                file=meta.get("file", os.path.basename(file_path)),
            ),
            file_path=usable_path,
            presentation=prs,
        )
        ctx.info.layouts_count = sum(len(m.slide_layouts) for m in prs.slide_masters)

        # 按 layouts_hint 归类 layouts
        hints = meta.get("layouts_hint", {})
        ctx.cover_layouts = self._resolve_layouts(prs, hints.get("cover", []))
        ctx.section_layouts = self._resolve_layouts(prs, hints.get("section", []))
        ctx.content_layouts = self._resolve_layouts(prs, hints.get("content", []))
        ctx.picture_layouts = self._resolve_layouts(
            prs, hints.get("picture") or hints.get("content", [])
        )
        ctx.summary_layouts = self._resolve_layouts(
            prs, hints.get("summary") or hints.get("section", [])
        )

        self._templates[tid] = ctx
        logger.info(
            "Loaded template '%s' from %s (cover=%d, section=%d, content=%d)",
            tid, usable_path,
            len(ctx.cover_layouts), len(ctx.section_layouts), len(ctx.content_layouts),
        )

    def _sanitize_template(self, file_path: str) -> str:
        """检测并修复源模板中的重复 slide layout / master 名字。

        问题背景：用户手工制作的 .pptx 在某些 PowerPoint 版本下允许同 master 下
        存在重名 layout（保存时不强制唯一），但 python-pptx 重新写入时会因
        zipfile 内部文件冲突产生警告并可能生成不合法 pptx。

        解决：首次加载时把模板去重后写入 cache 子目录，后续直接用 cache 版本。
        """
        cache_dir = os.path.join(TEMPLATES_DIR, "cache")
        # 缓存文件路径（按源文件 mtime 生成版本号，避免热更新失效）
        try:
            mtime = int(os.path.getmtime(file_path))
        except OSError:
            mtime = 0
        cached = os.path.join(cache_dir, f"{os.path.basename(file_path)}.{mtime}.clean.pptx")

        if os.path.exists(cached):
            return cached

        try:
            src_prs = Presentation(file_path)
        except Exception as e:
            logger.warning("_sanitize_template: cannot read %s: %s", file_path, e)
            return file_path

        # 检测重复：把所有 layout name 收集，统计重复
        seen_names: Dict[str, int] = {}
        for master in src_prs.slide_masters:
            for lay in master.slide_layouts:
                nm = lay.name or ""
                seen_names[nm] = seen_names.get(nm, 0) + 1
        has_dup_layout = any(c > 1 for c in seen_names.values())

        # master 重名
        seen_masters: Dict[str, int] = {}
        for master in src_prs.slide_masters:
            nm = master.name or f"master_{id(master)}"
            seen_masters[nm] = seen_masters.get(nm, 0) + 1
        has_dup_master = any(c > 1 for c in seen_masters.values())

        if not has_dup_layout and not has_dup_master:
            return file_path

        # 有重复：通过修改 partname 让它们在 OOXML 中唯一
        # 实际上 layout.name 只是显示名，不影响 partname。但部分 PowerPoint 工具
        # 会用 name 生成 partname，导致 zipfile 冲突。我们直接重命名重复 layout。
        for mi, master in enumerate(src_prs.slide_masters):
            used = set()
            for li, lay in enumerate(master.slide_layouts):
                base = lay.name or f"Layout{li}"
                new_name = base
                suffix = 2
                while new_name in used:
                    new_name = f"{base} ({suffix})"
                    suffix += 1
                used.add(new_name)
                if new_name != lay.name:
                    try:
                        lay.name = new_name
                    except Exception:
                        pass

        # master 同样处理
        used_m = set()
        for mi, master in enumerate(src_prs.slide_masters):
            base = master.name or f"Master{mi}"
            new_name = base
            suffix = 2
            while new_name in used_m:
                new_name = f"{base} ({suffix})"
                suffix += 1
            used_m.add(new_name)
            if new_name != master.name:
                try:
                    master.name = new_name
                except Exception:
                    pass

        os.makedirs(cache_dir, exist_ok=True)
        try:
            src_prs.save(cached)
            logger.info("_sanitize_template: sanitized %s -> %s", file_path, cached)
            return cached
        except Exception as e:
            logger.warning("_sanitize_template: save sanitized copy failed: %s; using original", e)
            return file_path

    def _resolve_layouts(self, prs: Presentation, names: List[str]) -> List[Any]:
        """按名称列表从所有 master 的 layouts 中匹配"""
        if not names:
            return []
        result = []
        all_layouts = []
        for m in prs.slide_masters:
            for lay in m.slide_layouts:
                all_layouts.append(lay)
        for name in names:
            for lay in all_layouts:
                if lay.name == name and lay not in result:
                    result.append(lay)
                    break
        return result

    # ── 内部：报告解析 ───────────────────────────────────────

    def _parse_outline(self, report: Report) -> Dict[str, Any]:
        """把报告 Markdown 拆为：标题 + 章节列表(每章节含 title/content/level)"""
        content = report.content_markdown or ""
        lines = content.split("\n")

        title = report.title
        sections: List[Dict[str, Any]] = []
        current: Optional[Dict[str, Any]] = None

        for raw in lines:
            line = raw.rstrip()
            if line.startswith("## "):
                # 切到下一章节
                if current is not None:
                    current["content"] = current["content"].strip()
                    sections.append(current)
                current = {"title": line[3:].strip(), "content": "", "level": 2}
            elif line.startswith("# ") and not sections and current is None:
                title = line[2:].strip()
            else:
                if current is None:
                    # ## 出现之前的内容，丢掉
                    continue
                current["content"] += line + "\n"

        if current is not None:
            current["content"] = current["content"].strip()
            sections.append(current)

        return {"title": title, "sections": sections}

    def _split_section_content(self, text: str) -> List[str]:
        """按段落 + 字符数双约束切分内容,保护 ![alt](data:image/...) 完整不被切碎。"""
        text = (text or "").strip()
        if not text:
            return [""]
        if len(text) <= MAX_CHARS_PER_CONTENT_SLIDE:
            return [text]

        # 按段落切分(双换行作为段落边界)
        paragraphs = re.split(r'\n\s*\n', text)
        # 同时把单换行的列表项作为独立块
        all_blocks: List[str] = []
        for p in paragraphs:
            for sub in re.split(r'\n', p):
                if sub.strip():
                    all_blocks.append(sub.strip())

        chunks: List[str] = []
        current = ""
        for block in all_blocks:
            tentative = (current + "\n" + block) if current else block
            if len(tentative) <= MAX_CHARS_PER_CONTENT_SLIDE:
                current = tentative
            else:
                if current:
                    chunks.append(current)
                # 单段过长,按句子再切
                if len(block) > MAX_CHARS_PER_PARAGRAPH:
                    sub_chunks = self._split_long_paragraph(block)
                    if len(sub_chunks) > 1:
                        chunks.extend(sub_chunks[:-1])
                        current = sub_chunks[-1]
                    else:
                        current = sub_chunks[0] if sub_chunks else ""
                else:
                    current = block
        if current:
            chunks.append(current)
        return chunks if chunks else [""]

    def _split_long_paragraph(self, paragraph: str) -> List[str]:
        """对超长段落按句子边界切分,保护图片语法不被切碎。"""
        if len(paragraph) <= MAX_CHARS_PER_PARAGRAPH:
            return [paragraph]

        chunks: List[str] = []
        remaining = paragraph
        while len(remaining) > MAX_CHARS_PER_PARAGRAPH:
            cut = -1
            for sep in ["。", "！", "？", "；", "\n", "，", "、", " "]:
                pos = remaining.rfind(sep, 0, MAX_CHARS_PER_PARAGRAPH)
                if pos > cut:
                    cut = pos + len(sep)
            # 检查 cut 是否在图片内部,如果是则跳到图片结束
            img_pattern = re.compile(r'!\[([^\]]*)\]\(data:image/[A-Za-z0-9+/=]+\)', re.DOTALL)
            for m in img_pattern.finditer(remaining):
                if m.start() < cut < m.end():
                    cut = m.end()
                    break
            if cut <= 0:
                cut = MAX_CHARS_PER_PARAGRAPH
            chunks.append(remaining[:cut].strip())
            remaining = remaining[cut:].strip()
        if remaining:
            chunks.append(remaining)
        return chunks

    # ── 内部：占位符填充 ─────────────────────────────────────

    def _find_placeholder(self, slide, type: Optional[int] = None, idx: Optional[int] = None, name_contains: Optional[str] = None):
        """查找占位符,按 type / idx / name 依次匹配"""
        for ph in slide.placeholders:
            try:
                if type is not None and ph.placeholder_format.type == type:
                    return ph
                if idx is not None and ph.placeholder_format.idx == idx:
                    return ph
                if name_contains and name_contains in (ph.name or ""):
                    return ph
            except Exception:
                continue
        return None

    def _set_placeholder_text(self, slide, idx_predicate, text: str, font_size: Optional[int] = None, bold: Optional[bool] = None, color: Optional[RGBColor] = None, alignment: Optional[int] = None):
        """找到第一个满足 idx_predicate(idx, type, name) 的占位符并填充文本"""
        for ph in slide.placeholders:
            if idx_predicate(ph.placeholder_format.idx, ph.placeholder_format.type, ph.name):
                tf = ph.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                # 默认左对齐,覆盖模板的居中
                p.alignment = alignment if alignment is not None else PP_ALIGN.LEFT
                run = p.add_run()
                run.text = text
                # 统一字体
                _set_run_font(run, size=font_size, bold=bold, color=color)
                return True
        return False

    def _add_cover_slide(self, prs: Presentation, ctx: TemplateContext, report: Report):
        layout = ctx.pick_cover(0)
        slide = prs.slides.add_slide(layout)
        # 标题占位符:优先 CENTER_TITLE(3) > TITLE(1)
        self._set_placeholder_text(
            slide,
            lambda i, t, n: t in (3, 1) or "标题" in n,
            report.title or "",
            font_size=32,
            bold=True,
            color=DEFAULT_TITLE_COLOR,
            alignment=PP_ALIGN.CENTER if ctx.pick_cover(0) else PP_ALIGN.LEFT,
        )
        # 副标题占位符
        subtitle = report.summary or datetime.now().strftime("%Y.%m")
        self._set_placeholder_text(
            slide,
            lambda i, t, n: t == 4 or "副标题" in n,
            subtitle,
            font_size=16,
            color=DEFAULT_TEXT_COLOR,
            alignment=PP_ALIGN.CENTER,
        )
        # 清空其他未填充的占位符,避免显示"单击此处添加文本"
        for ph in list(slide.placeholders):
            try:
                tf = ph.text_frame
                if tf.text and tf.text.strip() and tf.text not in ("单击此处添加文本", "Click to add text"):
                    continue
                tf.clear()
                p = tf.paragraphs[0]
                p.text = ""
            except Exception:
                pass

    def _add_toc_slide(self, prs: Presentation, ctx: TemplateContext, title: str, items: List[str]):
        layout = ctx.pick_content(0) if ctx.content_layouts else ctx.pick_cover(0)
        slide = prs.slides.add_slide(layout)
        # 标题
        self._set_placeholder_text(
            slide,
            lambda i, t, n: t in (1, 3) or "标题" in n,
            "目录",
            font_size=28,
            bold=True,
            color=DEFAULT_TITLE_COLOR,
        )
        # 优先使用母版 BODY/OBJECT 占位符填充目录项
        body_ph = self._find_placeholder(slide, type=7) or self._find_placeholder(slide, type=2)
        if body_ph is None:
            # 兜底: 用自建 textbox
            body_ph = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.0), Inches(5.5))
        # 调整 body 区域大小(避免占位符过小)
        try:
            body_ph.left = Inches(0.8)
            body_ph.top = Inches(1.6)
            body_ph.width = Inches(11.7)
            body_ph.height = Inches(5.5)
        except Exception:
            pass
        tf = body_ph.text_frame
        tf.word_wrap = True
        tf.clear()
        for i, item in enumerate(items, start=1):
            p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
            p.text = f"{i:02d}. {item}"
            p.level = 0
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(8)
            p.space_after = Pt(8)
            # 整段设置(实际 run 才会真正显示字体)
            run = p.runs[0] if p.runs else p.add_run()
            run.text = p.text
            _set_run_font(run, size=18, color=DEFAULT_TEXT_COLOR)

    def _add_section_cover_slide(self, prs: Presentation, ctx: TemplateContext, idx: int, section: Dict[str, Any]):
        layout = ctx.pick_section(idx - 1)
        slide = prs.slides.add_slide(layout)
        # 优先用母版 TITLE/CENTER_TITLE 占位符填充章节标题
        full_title = f"{idx:02d}  {section['title']}"
        title_set = self._set_placeholder_text(
            slide,
            lambda i, t, n: t in (1, 3) or "标题" in n,
            full_title,
            font_size=32,
            bold=True,
            color=DEFAULT_TITLE_COLOR,
        )
        if not title_set:
            # 兜底:在中央添加大字号章节标题
            tx = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(2.0))
            tf = tx.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = full_title
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.text = full_title
            _set_run_font(run, size=44, bold=True, color=DEFAULT_TITLE_COLOR)
        # 不再清空其他占位符 - 保留母版原样(背景/装饰)

    def _add_content_slide(self, prs: Presentation, ctx: TemplateContext, section_title: str, content: str, page_idx: int = 0):
        # 跳过空白内容,避免创建空幻灯片
        if self._is_content_empty(content):
            logger.debug(f"Skipping empty content slide for section '{section_title}'")
            return
        layout = ctx.pick_content(page_idx)
        slide = prs.slides.add_slide(layout)
        # 标题 - 优先使用母版 TITLE 占位符
        title_text = section_title if page_idx == 0 else f"{section_title} (续)"
        title_set = self._set_placeholder_text(
            slide,
            lambda i, t, n: t in (1, 3) or "标题" in n,
            title_text,
            font_size=22,
            bold=True,
            color=DEFAULT_TITLE_COLOR,
        )
        if not title_set:
            # 兜底:自建标题 textbox
            tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0))
            tf = tx.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = title_text
            p.alignment = PP_ALIGN.LEFT
            run = p.runs[0] if p.runs else p.add_run()
            run.text = title_text
            _set_run_font(run, size=24, bold=True, color=DEFAULT_TITLE_COLOR)

        # 检测内容是否含图片
        has_image = bool(re.search(r'!\[([^\]]*)\]\(data:image/[A-Za-z0-9+/=]+\)', content, re.DOTALL))

        # 内容区域:有图片用双区域版式,纯文字用单区域
        if has_image:
            # 上半文字(top=1.5, height=2.8 in),下半图片(top=4.5, height=2.7 in)
            self._render_content_with_images(slide, content, section_title)
        else:
            # 纯文字:使用母版 OBJECT/BODY 占位符
            self._render_text_only(slide, content)

    def _render_text_only(self, slide, content: str):
        """渲染纯文字内容到母版 OBJECT/BODY 占位符或自建 textbox。"""
        body_ph = self._find_placeholder(slide, type=7) or self._find_placeholder(slide, type=2)
        if body_ph is None:
            # 兜底:自建 textbox
            body_ph = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
        # 调整大小
        try:
            body_ph.left = Inches(0.5)
            body_ph.top = Inches(1.5)
            body_ph.width = Inches(12.3)
            body_ph.height = Inches(5.5)
        except Exception:
            pass
        # 提取纯文字(去除图片语法)
        text_only = re.sub(r'!\[([^\]]*)\]\(data:image/[A-Za-z0-9+/=]+\)', '', content, flags=re.DOTALL)
        text_only = text_only.strip()
        if not text_only:
            return
        tf = body_ph.text_frame
        tf.word_wrap = True
        tf.clear()
        # 自适应字号
        font_size = _calc_adaptive_font_size(len(text_only))
        self._fill_text_frame(tf, text_only, font_size=font_size)

    def _render_content_with_images(self, slide, content: str, section_title: str):
        """渲染"上文字下图片"双区域版式。"""
        # 上半文字区
        # 优先使用母版 OBJECT/BODY 占位符作为上半区域
        body_ph = self._find_placeholder(slide, type=7) or self._find_placeholder(slide, type=2)
        if body_ph is not None:
            try:
                body_ph.left = Inches(0.5)
                body_ph.top = Inches(1.5)
                body_ph.width = Inches(12.3)
                body_ph.height = Inches(2.8)
            except Exception:
                pass
            tf = body_ph.text_frame
            tf.word_wrap = True
            tf.clear()
        else:
            tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.8))
            tf = tx.text_frame
            tf.word_wrap = True
            tf.clear()

        # 填充文字(去除图片语法)
        text_only = re.sub(r'!\[([^\]]*)\]\(data:image/[A-Za-z0-9+/=]+\)', '', content, flags=re.DOTALL)
        text_only = text_only.strip()
        if text_only:
            font_size = _calc_adaptive_font_size(len(text_only))
            self._fill_text_frame(tf, text_only, font_size=font_size)

        # 下半图片区:从 content 中提取所有 base64 图片,插入到下半
        img_pattern = re.compile(r'!\[([^\]]*)\]\(data:image/png;base64,([A-Za-z0-9+/=\s]+)\)', re.DOTALL)
        matches = list(img_pattern.finditer(content))
        if not matches:
            return
        # 计算图片区(top=4.5, height=2.7 in)
        img_top = Inches(4.5)
        img_area_height = Inches(2.7)
        img_area_left = Inches(0.5)
        img_area_width = Inches(12.3)
        # 均匀分布多张图片
        n = len(matches)
        gap = Inches(0.2)
        each_width = Emu(int((img_area_width - gap * max(0, n - 1)) / n))
        for i, m in enumerate(matches):
            alt_text = m.group(1)
            b64_data = re.sub(r'\s+', '', m.group(2))
            try:
                img_blob = io.BytesIO(base64.b64decode(b64_data))
                # 验证图片有效性
                from PIL import Image as PILImage
                img_blob.seek(0)
                PILImage.open(img_blob).verify()
                img_blob.seek(0)
                left = img_area_left + Emu(int((each_width + gap) * i))
                # 限制最大高度
                from PIL import Image as PILImage
                img_blob.seek(0)
                with PILImage.open(img_blob) as pim:
                    w, h = pim.size
                aspect = h / w if w else 0.75
                # 等比缩放: 宽 each_width, 高等比
                max_h = Inches(2.5)
                pic_w = each_width
                pic_h = Emu(int(pic_w * aspect))
                if pic_h > max_h:
                    pic_h = max_h
                    pic_w = Emu(int(pic_h / aspect))
                # 垂直居中
                v_offset = Emu(int((img_area_height - pic_h) / 2))
                slide.shapes.add_picture(
                    img_blob, left, img_top + v_offset,
                    width=pic_w, height=pic_h,
                )
                logger.debug(f"Embedded image in slide '{section_title}': {alt_text[:30]}")
            except Exception as e:
                logger.warning(f"Failed to embed image '{alt_text}': {e}")

    def _is_content_empty(self, text: str) -> bool:
        """检测内容是否为空（仅包含空白字符、分隔线、图片语法等）"""
        if not text:
            return True
        # 移除图片语法后检查
        cleaned = re.sub(r'!\[.*?\]\(data:image/[^)]+\)', '', text, flags=re.DOTALL)
        cleaned = re.sub(r'---+', '', cleaned)
        cleaned = re.sub(r'\*\*数据点\*\*[：:]\s*\n(?:[-•·]\s*.+\n?)*', '', cleaned)
        return not cleaned.strip()

    def _fill_text_frame(self, tf, content: str, slide=None, font_size: Optional[int] = None):
        """把 Markdown 段落塞到 text_frame 中,每段一个 paragraph。
        统一使用 _set_run_font 设置中英文字体,确保视觉一致。
        """
        default_size = font_size or 14
        # 预处理:把图片语法单独提取出来(由 _render_content_with_images 处理图片)
        img_pattern = re.compile(r'!\[([^\]]*)\]\(data:image/png;base64,([A-Za-z0-9+/=\s]+)\)', re.DOTALL)

        # 将 content 按图片语法切分,保留图片语法作为特殊 token
        parts = []
        last_end = 0
        for m in img_pattern.finditer(content):
            if m.start() > last_end:
                parts.append(('text', content[last_end:m.start()]))
            parts.append(('image', m.group(1), m.group(2).replace('\n', '').replace('\r', '').replace(' ', '')))
            last_end = m.end()
        if last_end < len(content):
            parts.append(('text', content[last_end:]))
        if not parts:
            parts = [('text', '')]

        para_count = 0
        for part in parts:
            if part[0] == 'image':
                # 在 _render_content_with_images 中已统一处理图片,这里只放占位文字
                alt_text = part[1]
                p = tf.paragraphs[0] if para_count == 0 else tf.add_paragraph()
                p.text = f"[图片: {alt_text}]"
                p.alignment = PP_ALIGN.CENTER
                # 整段字体设置
                run = p.runs[0] if p.runs else p.add_run()
                run.text = p.text
                _set_run_font(run, size=11, color=RGBColor(0x66, 0x66, 0x66), italic=True)
                para_count += 1
                continue

            text_block = part[1].strip()
            if not text_block:
                continue
            for p_text in [p.strip() for p in re.split(r"\n+", text_block) if p.strip()]:
                p = tf.paragraphs[0] if para_count == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(4)
                # 处理 ### 三级标题
                if p_text.startswith("### "):
                    p_text = p_text[4:]
                    size = max(default_size + 1, 16)
                    color = DEFAULT_TITLE_COLOR
                    bold = True
                elif p_text.startswith("## "):
                    p_text = p_text[3:]
                    size = max(default_size + 3, 18)
                    color = DEFAULT_TITLE_COLOR
                    bold = True
                elif p_text.startswith("- ") or p_text.startswith("* "):
                    p_text = "• " + p_text[2:]
                    size = max(default_size - 1, 12)
                    color = DEFAULT_TEXT_COLOR
                    bold = False
                else:
                    size = default_size
                    color = DEFAULT_TEXT_COLOR
                    bold = False
                p.text = p_text
                # 整段字体设置(关键:用 _set_run_font 统一中英文字体)
                run = p.runs[0] if p.runs else p.add_run()
                run.text = p_text
                _set_run_font(run, size=size, bold=bold, color=color)
                para_count += 1

        if para_count == 0:
            # 至少添加一个段落占位
            p = tf.paragraphs[0]
            p.text = ""
            p.alignment = PP_ALIGN.LEFT
            run = p.runs[0] if p.runs else p.add_run()
            run.text = ""
            _set_run_font(run, size=default_size, color=DEFAULT_TEXT_COLOR)

    def _add_picture_slide(self, prs: Presentation, ctx: TemplateContext, chart: Dict[str, str]):
        layout = ctx.pick_picture(0)
        slide = prs.slides.add_slide(layout)
        title = chart.get("title", "图表")
        # 标题
        self._set_placeholder_text(
            slide,
            lambda i, t, n: t in (1, 3) or "标题" in n,
            title,
            font_size=20,
            bold=True,
        )
        # 图片占位符优先，否则直接 add_picture
        pic_added = False
        for ph in slide.placeholders:
            if ph.placeholder_format.type == 18:  # PICTURE
                try:
                    img_blob = io.BytesIO(base64.b64decode(chart["base64"]))
                    ph.insert_picture(img_blob)
                    pic_added = True
                except Exception as e:
                    logger.warning("insert_picture into placeholder failed: %s", e)
                break
        if not pic_added:
            try:
                img_blob = io.BytesIO(base64.b64decode(chart["base64"]))
                # 缩放至画布 85% 宽，居中
                slide_w = prs.slide_width
                slide_h = prs.slide_height
                max_w = int(slide_w * 0.85)
                max_h = int(slide_h * 0.65)
                from PIL import Image as PILImage
                try:
                    pil = PILImage.open(img_blob)
                    iw, ih = pil.size
                    scale = min(max_w / iw, max_h / ih, 1.0)
                    pic = slide.shapes.add_picture(
                        img_blob,
                        int((slide_w - iw * scale) / 2),
                        Inches(1.5),
                        width=int(iw * scale),
                        height=int(ih * scale),
                    )
                except Exception:
                    # 装不了 PIL 时的兜底
                    slide.shapes.add_picture(img_blob, Inches(1), Inches(1.5), width=max_w)
            except Exception as e:
                logger.warning("Failed to add picture to slide: %s", e)

    def _add_summary_slide(self, prs: Presentation, ctx: TemplateContext, report: Report, sections: List[Dict[str, Any]]):
        layout = ctx.pick_summary(0)
        slide = prs.slides.add_slide(layout)
        self._set_placeholder_text(
            slide,
            lambda i, t, n: t in (1, 3) or "标题" in n,
            "总结与展望",
            font_size=28,
            bold=True,
        )
        # 删除模板 BODY/OBJECT 占位符,使用自建 textbox
        placeholders_to_remove = []
        for ph in list(slide.placeholders):
            try:
                if ph.placeholder_format.type in (2, 7):
                    placeholders_to_remove.append(ph)
            except Exception:
                pass
        for ph in placeholders_to_remove:
            try:
                sp = ph._element
                sp.getparent().remove(sp)
            except Exception:
                pass

        last = sections[-1]["content"] if sections else (report.summary or "")
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9.0), Inches(5.5))
        tf = tx.text_frame
        tf.word_wrap = True
        tf.clear()
        if last:
            self._fill_text_frame(tf, last, slide=slide)
        else:
            p = tf.paragraphs[0]
            p.text = "本报告已完成全部维度的分析与总结"
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ── 内部：兜底 ───────────────────────────────────────────

    def _fallback_pptx(self, report: Report, chart_images: Optional[List[Dict[str, str]]] = None) -> bytes:
        """无任何模板可用时的最简版 PPT（与原 generate_pptx 行为一致）"""
        from pptx import Presentation as DefaultPresentation
        prs = DefaultPresentation()
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = report.title
        if len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = report.summary or ""
        content_slide = prs.slides.add_slide(prs.slide_layouts[1])
        content_slide.shapes.title.text = "报告内容"
        body = content_slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        tf.text = (report.content_markdown or "")[:2000]
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    # ── 内部：zip 去重 ───────────────────────────────────────

    def _dedupe_zip_parts(self, data: bytes) -> bytes:
        """修复 python-pptx 写出的 zip 内部重名 part。

        背景：用户提供的 .pptx 模板中多 master 场景下，python-pptx 重新打包时会
        出现多个相同 partname（如多个 slideLayout1.xml），导致 PowerPoint 打开
        时报"文件已损坏"。

        策略：按出现顺序给每个 entry 分配唯一名称，遇到重复时依次加 _dup1/_dup2/...
        并同步更新所有 .rels 中对旧 partname 的引用（含相对路径形式）。
        """
        import zipfile

        try:
            src_zip = zipfile.ZipFile(io.BytesIO(data), "r")
        except Exception:
            return data

        # 第一遍：按出现顺序去重
        names = src_zip.namelist()
        seen_count: Dict[str, int] = {}
        # key: 原 zip entry 在 names 列表中的索引；value: 最终写入的 entry 名
        index_to_name: Dict[int, str] = {}
        rename_map: Dict[str, str] = {}  # 原 partname -> 最终唯一名（最后一次出现的）

        for idx, n in enumerate(names):
            if n in seen_count:
                seen_count[n] += 1
                new_name = _make_unique(n, seen_count[n])
            else:
                seen_count[n] = 0
                new_name = n
            index_to_name[idx] = new_name
            # 记录"原 partname 最后一个被重命名后的名字"用于 .rels 替换
            if new_name != n:
                rename_map[n] = new_name

        if not rename_map:
            src_zip.close()
            return data

        # 第二遍：写新 zip
        out = io.BytesIO()
        with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as dst_zip:
            for idx, n in enumerate(names):
                content = src_zip.read(n)
                dst_zip.writestr(index_to_name[idx], content)

        # 第三遍：修复 .rels 文件中对旧 partname 的引用
        # 注意：rels 中 Target 是相对路径（如 slideLayouts/slideLayout1.xml），
        # 因此替换时要兼容"完整路径"和"裸文件名"两种形式。
        src_zip.close()
        out.seek(0)
        fixed = io.BytesIO()
        with zipfile.ZipFile(out, "r") as in_z, zipfile.ZipFile(fixed, "w", zipfile.ZIP_DEFLATED) as out_z:
            for item in in_z.infolist():
                content = in_z.read(item.filename)
                if item.filename.endswith(".rels"):
                    text = content.decode("utf-8", errors="ignore")
                    for old, new in rename_map.items():
                        # 完整路径形式
                        text = text.replace(f'"{old}"', f'"{new}"')
                        text = text.replace(f'/{old}"', f'/{new}"')
                        # 裸文件名形式（rels 中常省略父目录）
                        bare = old.rsplit("/", 1)[-1]
                        new_bare = new.rsplit("/", 1)[-1]
                        if bare != old:
                            text = text.replace(f'"{bare}"', f'"{new_bare}"')
                    content = text.encode("utf-8")
                out_z.writestr(item, content)

        return fixed.getvalue()

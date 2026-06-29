"""
HTML → PPT 转换器（兼容接口 + 原生渲染委托）

v2: 委托到 PPTXNativeRenderer 做原生 python-pptx 渲染（可编辑文字/图片），
    不再使用截图方案。

旧版行为（v1）：Playwright 截图 → 插入 PPT 位图
新版行为（v2）：HTMLPageModel → python-pptx 原生元素
"""

import logging
import os
from typing import List, Optional

from app.services.html_report_generator import HTMLPageModel
from app.services.pptx_native_renderer import PPTXNativeRenderer

logger = logging.getLogger(__name__)


class HTMLToPPTConverter:
    """
    HTML → PPT 转换器（兼容接口）

    v2: 委托到 PPTXNativeRenderer 做原生渲染。
        接受 HTML 字符串，内部解析 slide 列表，使用原生元素生成 PPTX。

    如果调用方有 HTMLPageModel 列表，直接使用 PPTXNativeRenderer.render() 更快。
    """

    def __init__(self):
        pass

    async def convert(
        self,
        html: str,
        output_path: str,
        template_id: Optional[str] = None,
        pages: Optional[List[HTMLPageModel]] = None,
    ):
        """将 HTML/PageModel 转换为 PPT

        Args:
            html: HTML 字符串（当 pages=None 时使用，降级到旧版截图）
            output_path: PPT 输出路径
            template_id: PPT 模板 ID（可选）
            pages: HTML 页面模型列表（优先使用，原生渲染）
        """
        template_path = self._resolve_template_path(template_id)

        if pages:
            # 新版：原生渲染
            logger.info(f"PPTXNativeRenderer: rendering {len(pages)} slides natively")
            renderer = PPTXNativeRenderer(template_path=template_path)
            renderer.render(pages, output_path, title="Market Insight Report")
            logger.info(f"Native PPTX saved: {output_path}")
        else:
            # 降级：截图方案（当没有 PageModel 时）
            logger.warning(
                "HTMLToPPTConverter: no pages provided, falling back to screenshot mode. "
                "For best quality, pass HTMLPageModel list."
            )
            await self._convert_screenshot_fallback(html, output_path, template_path)

    async def _convert_screenshot_fallback(
        self, html: str, output_path: str, template_path: Optional[str] = None,
    ):
        """降级到截图方案（旧行为）"""
        import tempfile
        from pptx import Presentation
        from app.services.playwright_renderer import PlaywrightRenderer

        renderer = PlaywrightRenderer()
        with tempfile.TemporaryDirectory() as tmpdir:
            screenshot_paths = await renderer.render_to_screenshots(html, tmpdir)
            if not screenshot_paths:
                raise ValueError("No screenshots generated from HTML")

            prs = Presentation(template_path) if template_path else Presentation()
            for sp in screenshot_paths:
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.add_picture(
                    sp, left=0, top=0,
                    width=prs.slide_width, height=prs.slide_height,
                )
            prs.save(output_path)
            logger.info(f"Fallback PPT saved: {output_path}")

    def _resolve_template_path(self, template_id: Optional[str]) -> Optional[str]:
        """根据 template_id 解析模板文件路径"""
        import glob
        
        # 优先查 PPT 模板缓存目录
        cache_dir = os.path.join(
            os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
            "templates", "ppt", "cache"
        )
        if os.path.isdir(cache_dir):
            # 查找匹配的模板文件 (template{id}.pptx.*.clean.pptx)
            pattern = f"template{template_id}.pptx.*.clean.pptx" if template_id else "template1.pptx.*.clean.pptx"
            matches = glob.glob(os.path.join(cache_dir, pattern))
            if matches:
                # 取最新的
                path = max(matches, key=os.path.getmtime)
                logger.info(f"Resolved template path from cache: {path}")
                return path
        
        # 其次查工作目录
        if template_id:
            candidates = [
                f"模板{template_id}.pptx",
                f"template_{template_id}.pptx",
                os.path.join("backend", "data", "pptx", f"template_{template_id}.pptx"),
            ]
            for path in candidates:
                if os.path.exists(path):
                    return path
        
        # 无 template_id：尝试默认模板1
        if not template_id:
            default = self._resolve_template_path("1")
            if default:
                return default
        
        logger.warning("No PPT template found, using default blank presentation")
        return None
